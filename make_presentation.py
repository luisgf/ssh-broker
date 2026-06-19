"""
SSH Broker – Corporate presentation generator (python-pptx)
Style: Zara-inspired editorial — Didone serif display + minimal sans, monochrome
Slides: 36
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x0d, 0x0d, 0x0d)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY1  = RGBColor(0x1A, 0x1A, 0x1A)
GRAY2  = RGBColor(0x3A, 0x3A, 0x3A)
GRAY3  = RGBColor(0x6B, 0x6B, 0x6B)
GRAY4  = RGBColor(0xB0, 0xB0, 0xB0)
GRAY5  = RGBColor(0xE8, 0xE8, 0xE8)
GRAY6  = RGBColor(0xF5, 0xF5, 0xF5)

# ── Typography (Zara-inspired editorial) ─────────────────────────────────────
# Display: high-contrast Didone serif (Zara's Didot logo). Body / labels: minimal
# grotesque sans with wide tracking. Code: monospace. Names target macOS; a
# missing font falls back to the renderer default.
SERIF = "Didot"
SANS  = "Helvetica Neue"
MONO  = "Menlo"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

prs.core_properties.author           = "Luis Gonzalez Fernandez"
prs.core_properties.last_modified_by = "Luis Gonzalez Fernandez"
prs.core_properties.title            = "SSH Broker — Secure SSH Access for AI Agents"

# ── Helpers ────────────────────────────────────────────────────────────────

def solid_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def no_fill(shape):
    shape.fill.background()

def no_line(shape):
    shape.line.fill.background()

def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        solid_fill(shape, fill_color)
    else:
        no_fill(shape)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        no_line(shape)
    return shape

def _set_tracking(run, tracking):
    # tracking in points → OOXML spc attribute (units of 1/100 pt)
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(tracking * 100)))

def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(12), bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True,
                serif=None, tracking=None, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # Zara typography: Didone serif for display sizes, minimal sans for body.
    if font_name is None:
        use_serif = serif if serif is not None else (font_size.pt >= 24)
        font_name = SERIF if use_serif else SANS
    run.font.name = font_name
    if tracking is not None:
        _set_tracking(run, tracking)
    return txBox

def add_label_in_rect(slide, text, left, top, width, height,
                       fill_color=BLACK, text_color=WHITE,
                       font_size=Pt(11), bold=True):
    r = add_rect(slide, left, top, width, height, fill_color=fill_color)
    tf = r.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = SANS
    _set_tracking(run, 1.5)
    r.text_frame._txBody.attrib['anchor'] = 'ctr'
    return r

def add_arrow(slide, x1, y1, x2, y2, color=GRAY3, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    ln = connector.line._ln
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'none')
    headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'arrow')
    headEnd.set('w', 'med')
    headEnd.set('len', 'med')
    return connector

def slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def top_bar(slide, bg=BLACK, height=Inches(0.08)):
    add_rect(slide, 0, 0, SLIDE_W, height, fill_color=bg)

def bottom_bar(slide, text="SSH Broker · Confidential", bg=GRAY1, height=Inches(0.45)):
    add_rect(slide, 0, SLIDE_H - height, SLIDE_W, height, fill_color=bg)
    add_textbox(slide, text,
                Inches(0.35), SLIDE_H - height + Inches(0.08),
                Inches(10), height,
                font_size=Pt(9), color=GRAY4, align=PP_ALIGN.LEFT)

def slide_number(slide):
    # Number is derived from slide position so inserting slides never
    # requires renumbering call sites. The cover (slide 1) omits the call.
    n = len(prs.slides)
    add_textbox(slide, str(n),
                SLIDE_W - Inches(0.7), SLIDE_H - Inches(0.43),
                Inches(0.5), Inches(0.38),
                font_size=Pt(9), color=GRAY4, align=PP_ALIGN.RIGHT)

def section_label(slide, text, y=Inches(1.15)):
    add_textbox(slide, text.upper(),
                Inches(0.9), y, Inches(11), Inches(0.35),
                font_size=Pt(9), bold=True, color=GRAY4, tracking=2.5)

def title_text(slide, text, y=Inches(1.55), size=Pt(38)):
    add_textbox(slide, text,
                Inches(0.9), y, Inches(11.4), Inches(1.6),
                font_size=size, bold=True, color=BLACK, serif=True)

def mono_block(slide, text, left, top, width, height, bg=GRAY1, fg=GRAY4):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_textbox(slide, text, left + Inches(0.2), top + Inches(0.15),
                width - Inches(0.3), height - Inches(0.25),
                font_size=Pt(9), color=fg, wrap=True, font_name=MONO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)

add_rect(s, 0, Inches(3.3), Inches(0.18), Inches(1.9), fill_color=WHITE)

add_textbox(s, "SSH BROKER",
            Inches(0.55), Inches(1.5), Inches(11), Inches(1.4),
            font_size=Pt(64), bold=True, color=WHITE)

add_textbox(s, "Secure SSH access infrastructure\nfor AI agents and human operators",
            Inches(0.55), Inches(3.0), Inches(9), Inches(1.4),
            font_size=Pt(20), color=GRAY4)

add_textbox(s, "Ephemeral credentials · Zero static keys · Cryptographic audit trail",
            Inches(0.55), Inches(4.6), Inches(9), Inches(0.6),
            font_size=Pt(13), color=GRAY3, italic=True)

add_textbox(s, "June 2026  ·  v1.18.0",
            Inches(0.55), Inches(6.1), Inches(5), Inches(0.5),
            font_size=Pt(11), color=GRAY3)

add_textbox(s, "Luis Gonzalez Fernandez",
            Inches(0.55), Inches(6.6), Inches(7), Inches(0.4),
            font_size=Pt(10), color=GRAY3, italic=True)

# Minimal editorial cover — negative space replaces the decorative grid.

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "AGENDA",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)

add_textbox(s, "What we will cover.",
            Inches(0.9), Inches(1.38), Inches(9), Inches(0.75),
            font_size=Pt(34), bold=True, color=WHITE)

sections = [
    ("01", "Problem & Threat Model",
     "Why static SSH keys are unacceptable for AI agents · Attack vectors · What we defend against"),
    ("02", "Architecture & Core Mechanism",
     "How the broker works · Component map · Ephemeral certificates · MCP tools · Sudo elevation"),
    ("03", "Security Controls",
     "AI-action firewall (command policy + shell AST parsing, approval gate, behaviour guardrails) · Teams notifications · RBAC · Audit trail"),
    ("04", "Deployment",
     "Local mode (stdio) · Remote mode (HTTP + OAuth 2.1) · Microsoft Entra ID integration"),
    ("05", "Operations & Roadmap",
     "Day-to-day operations · Gaps toward production · Competitive landscape · Next steps"),
]

for i, (num, title, desc) in enumerate(sections):
    y = Inches(2.4) + i * Inches(0.95)
    add_textbox(s, num,
                Inches(0.9), y, Inches(0.7), Inches(0.8),
                font_size=Pt(28), bold=True, color=GRAY3)
    add_textbox(s, title,
                Inches(1.65), y, Inches(11), Inches(0.38),
                font_size=Pt(14), bold=True, color=WHITE)
    add_textbox(s, desc,
                Inches(1.65), y + Inches(0.4), Inches(11), Inches(0.45),
                font_size=Pt(10), color=GRAY3)
    add_rect(s, Inches(0.9), y + Inches(0.88), Inches(11.5), Inches(0.02),
             fill_color=GRAY2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "01  Problem & Threat Model")
title_text(s, "AI agents need to run\ncommands on Linux servers.", y=Inches(1.55), size=Pt(34))

cards = [
    ("Static SSH keys are\nexfiltrable",
     "An AI agent holding a long-lived SSH private key can be tricked via prompt injection or memory dump. Once stolen, the key works forever."),
    ("No audit trail\nfor AI actions",
     "Traditional SSH logs show the human account, not which AI model issued the command or what it was trying to achieve."),
    ("All-or-nothing\naccess control",
     "A single compromised key grants full access to every host it was configured for — with no per-command policy enforcement."),
]

for i, (title, body) in enumerate(cards):
    left = Inches(0.45) + i * Inches(4.2)
    add_rect(s, left, Inches(3.6), Inches(3.9), Inches(3.0),
             fill_color=WHITE, line_color=GRAY4, line_width=Pt(0.75))
    add_textbox(s, str(i + 1),
                left + Inches(0.2), Inches(3.75), Inches(0.45), Inches(0.55),
                font_size=Pt(28), bold=True, color=GRAY4)
    add_textbox(s, title,
                left + Inches(0.2), Inches(4.25), Inches(3.5), Inches(0.8),
                font_size=Pt(13), bold=True, color=BLACK)
    add_textbox(s, body,
                left + Inches(0.2), Inches(5.0), Inches(3.5), Inches(1.4),
                font_size=Pt(10.5), color=GRAY2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THREAT MODEL
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "01  Problem & Threat Model")
title_text(s, "Threat model", y=Inches(1.45), size=Pt(34))

add_textbox(s, "Five attack vectors — and the layer that stops each one.",
            Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.4),
            font_size=Pt(12), color=GRAY3, italic=True)

# Table header
col_w = [Inches(3.1), Inches(3.8), Inches(5.2)]
col_x = [Inches(0.5), Inches(3.6), Inches(7.4)]
hdrs  = ["Attack vector", "Without broker", "What stops it"]
TABLE_TOP = Inches(2.72)
ROW_H = Inches(0.74)

for j, (h, w, x) in enumerate(zip(hdrs, col_w, col_x)):
    add_label_in_rect(s, h, x, TABLE_TOP, w, Inches(0.42),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)

rows = [
    ("Prompt injection\n→ key exfiltration",
     "Attacker tricks the AI into leaking the static private key. Key is reusable forever.",
     "No static key exists. Ed25519 pair is generated in RAM per-operation and never written to disk."),
    ("Memory dump\nof broker process",
     "Attacker dumps process memory and extracts the long-lived SSH key.",
     "Only the ephemeral key lives in RAM during the operation. Certificate expires in 60–120 s."),
    ("Broker fully\ncompromised",
     "Attacker controls the broker binary and can do anything the key allows.",
     "Signer is isolated. A compromised broker cannot exceed the policy, elevate privileges, or forge certs."),
    ("AI agent\nbehaves anomalously",
     "Agent issues unusual commands or accesses hosts it has never touched. No one notices.",
     "Control plane behaviour guardrails detect rate spikes, new hosts, and novel commands. Escalates to approval."),
    ("High-risk command\nruns unsupervised",
     "A sensitive or destructive command executes with no human review and no real-time alert to operators.",
     "Approval gate holds the certificate until a human signs off — the request is surfaced in real time via a Microsoft Teams notification."),
]

for i, (vec, without, stops) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else GRAY5
    y  = TABLE_TOP + Inches(0.42) + i * ROW_H
    for j, (cell, w, x) in enumerate(zip([vec, without, stops], col_w, col_x)):
        add_rect(s, x, y, w, ROW_H, fill_color=bg,
                 line_color=GRAY4, line_width=Pt(0.5))
        fc = GRAY2 if j < 2 else BLACK
        add_textbox(s, cell, x + Inches(0.12), y + Inches(0.1),
                    w - Inches(0.2), ROW_H - Inches(0.15),
                    font_size=Pt(9.5), color=fc, bold=(j == 2))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "02  ARCHITECTURE & CORE MECHANISM",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)

add_textbox(s, "A broker that stands between\nthe AI and your infrastructure.",
            Inches(0.9), Inches(1.4), Inches(11), Inches(1.5),
            font_size=Pt(36), bold=True, color=WHITE)

points = [
    ("The AI model never touches a credential.",
     "It calls a tool. The broker handles everything else."),
    ("Ephemeral Ed25519 keys — generated in memory, used once, discarded.",
     "A stolen snapshot of the broker process contains no reusable secret."),
    ("Every action is signed, chained, and audited.",
     "Tamper-evident log correlates broker execution ↔ signer certificate ↔ sshd session."),
    ("Policy lives in the signer — isolated from the broker.",
     "Even a fully compromised broker cannot exceed what the signer allows."),
]

for i, (bold_part, light_part) in enumerate(points):
    y = Inches(3.05) + i * Inches(0.9)
    add_rect(s, Inches(0.85), y + Inches(0.12), Inches(0.08), Inches(0.08),
             fill_color=WHITE)
    tb = s.shapes.add_textbox(Inches(1.05), y, Inches(11.8), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = bold_part + " "
    r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = WHITE; r1.font.name = SANS
    r2 = p.add_run()
    r2.text = light_part
    r2.font.size = Pt(13); r2.font.color.rgb = GRAY4; r2.font.name = SANS

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "02  Architecture & Core Mechanism")
title_text(s, "How it works", y=Inches(1.5), size=Pt(32))

ROW1_Y = Inches(2.9)
BOX_H  = Inches(0.85)
BOX_W  = Inches(2.1)

boxes_row1 = [
    (Inches(0.5),  "AI Model\n(Claude / OpenCode)", GRAY1, WHITE),
    (Inches(3.15), "MCP Broker",                    GRAY2, WHITE),
    (Inches(5.8),  "Control Plane\n(PEP / guardrails)", GRAY3, WHITE),
    (Inches(8.45), "Signer\n(CA custodian)",         BLACK, WHITE),
]

for left, label, fc, tc in boxes_row1:
    add_label_in_rect(s, label, left, ROW1_Y, BOX_W, BOX_H,
                      fill_color=fc, text_color=tc, font_size=Pt(11))

arrow_labels = ["MCP tool call", "mTLS HTTPS", "mTLS HTTPS"]
for i in range(len(boxes_row1) - 1):
    x1 = boxes_row1[i][0] + BOX_W
    x2 = boxes_row1[i + 1][0]
    my = ROW1_Y + BOX_H / 2
    add_arrow(s, x1, my, x2, my, color=GRAY3, width=Pt(2))
    add_textbox(s, arrow_labels[i],
                x1 + Inches(0.05), ROW1_Y - Inches(0.4),
                x2 - x1 - Inches(0.05), Inches(0.35),
                font_size=Pt(8), color=GRAY3, align=PP_ALIGN.CENTER)

ROW2_Y = Inches(5.0)
targets = [(Inches(3.0), "Bastion :22"), (Inches(5.65), "Target host :22")]
for left, lbl in targets:
    add_label_in_rect(s, lbl, left, ROW2_Y, Inches(2.0), Inches(0.75),
                      fill_color=GRAY5, text_color=GRAY2, font_size=Pt(11), bold=False)
    add_rect(s, left, ROW2_Y, Inches(2.0), Inches(0.75),
             line_color=GRAY4, line_width=Pt(0.75))

mid_bx = boxes_row1[1][0] + BOX_W / 2
add_arrow(s, mid_bx, ROW1_Y + BOX_H, mid_bx, ROW2_Y, color=GRAY3, width=Pt(1.5))
add_textbox(s, "SSH + ephemeral cert",
            mid_bx - Inches(0.9), ROW1_Y + BOX_H + Inches(0.12),
            Inches(1.8), Inches(0.35), font_size=Pt(8), color=GRAY3)

add_arrow(s, targets[0][0] + Inches(2.0), ROW2_Y + Inches(0.375),
             targets[1][0],               ROW2_Y + Inches(0.375),
             color=GRAY3, width=Pt(1.5))
add_textbox(s, "ProxyJump",
            targets[0][0] + Inches(2.05), ROW2_Y - Inches(0.3),
            Inches(1.5), Inches(0.3), font_size=Pt(8), color=GRAY3)

add_textbox(s, "stdout / stderr / exit_code  →  AI model",
            Inches(0.5), Inches(6.1), Inches(12), Inches(0.4),
            font_size=Pt(10), color=GRAY3, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SYSTEM ARCHITECTURE / COMPONENTS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "02  ARCHITECTURE & CORE MECHANISM",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "System architecture — components",
            Inches(0.9), Inches(1.38), Inches(11.5), Inches(0.7),
            font_size=Pt(30), bold=True, color=WHITE)
add_textbox(s, "Four binaries · three trust zones. Policy and the CA key live only in the signer.",
            Inches(0.9), Inches(2.15), Inches(11.5), Inches(0.4),
            font_size=Pt(11), color=GRAY4, italic=True)

arch_zones = [
    ("CLIENT",
     ["AI Model",
      "mcp-broker  (stdio)",
      "mcp-broker-http  (HTTP + OIDC)"]),
    ("POLICY ENFORCEMENT\ncontrol-plane (PEP)",
     ["Approval Registry",
      "Behaviour Tracker",
      "Notifier:  Log / Webhook / Teams"]),
    ("CUSTODY\nsigner (PDP + CA custodian)",
     ["Policy Table",
      "Command Policy",
      "RBAC  (groups / user)",
      "CA private key"]),
]
panel_xs = [Inches(0.5), Inches(4.7), Inches(8.9)]
PW = Inches(3.9)
PY = Inches(2.7)
PH = Inches(3.0)
for px, (ztitle, boxes) in zip(panel_xs, arch_zones):
    add_rect(s, px, PY, PW, PH, fill_color=GRAY1, line_color=GRAY3, line_width=Pt(0.75))
    add_textbox(s, ztitle, px + Inches(0.2), PY + Inches(0.12),
                PW - Inches(0.4), Inches(0.55),
                font_size=Pt(10), bold=True, color=WHITE)
    for k, label in enumerate(boxes):
        by = PY + Inches(0.72) + k * Inches(0.55)
        is_ca = (label == "CA private key")
        add_label_in_rect(s, label, px + Inches(0.2), by,
                          PW - Inches(0.4), Inches(0.46),
                          fill_color=WHITE if is_ca else GRAY2,
                          text_color=BLACK if is_ca else WHITE,
                          font_size=Pt(9.5), bold=is_ca)

for i in range(len(panel_xs) - 1):
    x1 = panel_xs[i] + PW
    x2 = panel_xs[i + 1]
    my = PY + PH / 2
    add_arrow(s, x1, my, x2, my, color=GRAY4, width=Pt(1.5))
    add_textbox(s, "mTLS", x1, my - Inches(0.32),
                x2 - x1, Inches(0.28),
                font_size=Pt(7.5), color=GRAY4, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.5), fill_color=GRAY2)
add_textbox(s, "broker-ctl — management CLI:   host add / list   ·   approval allow / deny   ·   audit verify   ·   reload",
            Inches(0.65), Inches(6.03), Inches(12.0), Inches(0.34),
            font_size=Pt(9.5), bold=True, color=WHITE)
add_textbox(s, "Cross-cutting:   Audit logs ×3 (Ed25519-chained)   ·   PKI / mTLS   ·   SSH targets via ephemeral cert (bastion + hosts)",
            Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.34),
            font_size=Pt(9), color=GRAY4, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHY EPHEMERAL CERTIFICATES (ANATOMY)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "02  Architecture & Core Mechanism")
title_text(s, "Why ephemeral certificates, not keys", y=Inches(1.4), size=Pt(28))
add_textbox(s, "A certificate is a public key the CA signs with tight, short-lived constraints — there is no reusable secret to steal.",
            Inches(0.9), Inches(2.15), Inches(11.8), Inches(0.4),
            font_size=Pt(11), color=GRAY3, italic=True)

# Left: static key vs ephemeral cert
left_cards = [
    (GRAY3, "STATIC SSH KEY",
     "A reusable secret stored on disk or in the agent. Stolen once, it works forever — on every host it was configured for."),
    (BLACK, "EPHEMERAL CERTIFICATE",
     "A public key signed by the CA, valid 60–120 s and scoped to one host and command. Stolen, it has already expired and is narrowly scoped."),
]
for i, (hc, htxt, body) in enumerate(left_cards):
    cy = Inches(2.7) + i * Inches(1.95)
    add_rect(s, Inches(0.5), cy, Inches(5.7), Inches(1.75),
             fill_color=WHITE, line_color=GRAY4, line_width=Pt(0.75))
    add_label_in_rect(s, htxt, Inches(0.5), cy, Inches(5.7), Inches(0.45),
                      fill_color=hc, text_color=WHITE, font_size=Pt(11), bold=True)
    add_textbox(s, body, Inches(0.7), cy + Inches(0.6),
                Inches(5.35), Inches(1.05),
                font_size=Pt(10.5), color=GRAY2)

# Right: certificate anatomy
ax, ay, aw, ah = Inches(6.6), Inches(2.7), Inches(6.25), Inches(3.95)
add_rect(s, ax, ay, aw, ah, fill_color=WHITE, line_color=GRAY4, line_width=Pt(0.75))
add_label_in_rect(s, "SSH CERTIFICATE — what the CA signs",
                  ax, ay, aw, Inches(0.5),
                  fill_color=BLACK, text_color=WHITE, font_size=Pt(11), bold=True)
cert_fields = [
    ("Public key",       "Ed25519 — generated in RAM"),
    ("Principal",        "host:web01"),
    ("Validity",         "60–120 s  (short TTL)"),
    ("Critical options", "force-command · source-address · permit-pty"),
    ("Serial",           "1042  (audit correlation)"),
    ("Signed by",        "CA private key  (signer only)"),
]
for k, (label, val) in enumerate(cert_fields):
    fy = ay + Inches(0.65) + k * Inches(0.46)
    add_textbox(s, label, ax + Inches(0.2), fy,
                Inches(1.85), Inches(0.4), font_size=Pt(9.5), bold=True, color=GRAY2)
    add_textbox(s, val, ax + Inches(2.05), fy,
                aw - Inches(2.25), Inches(0.4), font_size=Pt(9.5), color=BLACK)
add_rect(s, ax + Inches(0.15), ay + ah - Inches(0.5), aw - Inches(0.3), Inches(0.4),
         fill_color=GRAY5)
add_textbox(s, "The private key never leaves the broker's RAM — only the public key is sent to be signed.",
            ax + Inches(0.25), ay + ah - Inches(0.46), aw - Inches(0.5), Inches(0.34),
            font_size=Pt(8.5), color=GRAY2, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EPHEMERAL CREDENTIALS STEP BY STEP
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "02  Architecture & Core Mechanism")
title_text(s, "Ephemeral credentials — step by step", y=Inches(1.5), size=Pt(30))

steps = [
    ("1", "AI calls ssh_execute(server, command)",
     "Tool call over stdio (local) or HTTPS bearer token (remote)."),
    ("2", "Broker generates Ed25519 key pair in RAM",
     "The private key never touches disk. There is nothing to steal."),
    ("3", "Broker sends Intent to Signer (mTLS)",
     "Intent contains: host, command, public key — never the private key."),
    ("4", "Signer checks policy and signs the certificate",
     "Principal, force-command, TTL, source-address — all baked into the cert."),
    ("5", "Broker dials SSH using the ephemeral cert",
     "Cert may be valid for as little as 60 seconds. After execution, it expires."),
    ("6", "Result returned; key material discarded",
     "stdout / stderr / exit_code reach the AI. No credentials linger."),
]

for i, (num, title, sub) in enumerate(steps):
    row = i // 3
    col = i % 3
    left = Inches(0.5) + col * Inches(4.2)
    top  = Inches(3.2) + row * Inches(1.7)
    add_rect(s, left, top, Inches(3.9), Inches(1.5),
             fill_color=WHITE, line_color=GRAY4, line_width=Pt(0.75))
    add_label_in_rect(s, num,
                      left + Inches(0.15), top + Inches(0.15),
                      Inches(0.38), Inches(0.38),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(11), bold=True)
    add_textbox(s, title,
                left + Inches(0.65), top + Inches(0.12),
                Inches(3.1), Inches(0.5),
                font_size=Pt(11), bold=True, color=BLACK)
    add_textbox(s, sub,
                left + Inches(0.15), top + Inches(0.7),
                Inches(3.6), Inches(0.7),
                font_size=Pt(10), color=GRAY2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THE 5 MCP TOOLS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "02  Architecture & Core Mechanism")
title_text(s, "What the AI can do —\nthe 5 MCP tools", y=Inches(1.4), size=Pt(32))

tool_headers = ["Tool", "Key parameters", "What it does"]
tool_rows = [
    ["ssh_list_servers",
     "—",
     "Returns all reachable hosts with their capabilities (sudo, PTY, jump). Always call first."],
    ["ssh_execute",
     "server, command\n[sudo, sudo_user, pty, ttl_seconds, dry_run]",
     "One-shot execution. Certificate bakes force-command. Result: stdout / stderr / exit_code."],
    ["ssh_session_open",
     "server\n[mode: exec | shell | pty, sudo, sudo_user, ttl_seconds]",
     "Opens a persistent SSH session. Reuses the connection for multiple commands."],
    ["ssh_session_exec",
     "session_id, command",
     "Runs a command inside an open session. No new certificate needed."],
    ["ssh_session_close",
     "session_id",
     "Closes the session and releases all resources. Certificate expires naturally after TTL."],
]

tcol_w = [Inches(2.4), Inches(3.2), Inches(6.7)]
tcol_x = [Inches(0.5), Inches(2.9), Inches(6.1)]
TABLE_TOP = Inches(2.95)
T_ROW_H   = Inches(0.72)

for j, (h, w, x) in enumerate(zip(tool_headers, tcol_w, tcol_x)):
    add_label_in_rect(s, h, x, TABLE_TOP, w, Inches(0.4),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)

for i, row in enumerate(tool_rows):
    bg = WHITE if i % 2 == 0 else GRAY5
    y  = TABLE_TOP + Inches(0.4) + i * T_ROW_H
    for j, (cell, w, x) in enumerate(zip(row, tcol_w, tcol_x)):
        add_rect(s, x, y, w, T_ROW_H, fill_color=bg,
                 line_color=GRAY4, line_width=Pt(0.5))
        is_tool = (j == 0)
        add_textbox(s, cell, x + Inches(0.1), y + Inches(0.08),
                    w - Inches(0.15), T_ROW_H - Inches(0.12),
                    font_size=Pt(9.5 if not is_tool else 10),
                    color=BLACK if is_tool else GRAY2,
                    bold=is_tool)

add_textbox(s,
            "Recommended flow:  ssh_list_servers  →  ssh_execute  (or ssh_session_open / exec / close)",
            Inches(0.5), Inches(6.7), Inches(12.5), Inches(0.35),
            font_size=Pt(10), color=GRAY3, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SUDO ELEVATION POLICY-GATED
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "02  ARCHITECTURE & CORE MECHANISM",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Sudo elevation — policy-gated\nby the signer, not the broker.",
            Inches(0.9), Inches(1.38), Inches(11.5), Inches(1.1),
            font_size=Pt(28), bold=True, color=WHITE)

# Two columns
for col, (mode, steps_list) in enumerate([
    ("One-shot  (ssh_execute + sudo=true)",
     [("1  Intent",        'broker → Intent{ sudo=true, command="id" }'),
      ("2  Policy check",  "signer: allow_sudo? allowed_sudo_users?"),
      ("3  Sign cert",     'force-command = "sudo -n -- /bin/sh -c \'id\'"'),
      ("4  SSH connect",   "sshd enforces force-command from the cert"),
      ("5  Result",        "stdout / stderr / exit_code → AI model"),
     ]),
    ("Session  (ssh_session_open + sudo=true)",
     [("1  Intent",        "broker → Intent{ sudo=true, purpose=session }"),
      ("2  Policy check",  "signer: same allow_sudo checks"),
      ("3  Sign cert",     'cert has no force-command'),
      ("4  ElevationPrefix", 'signer returns ElevationPrefix="sudo -n"'),
      ("5  Each exec",     'effective cmd: "sudo -n -- /bin/sh -c \'<cmd>\'"'),
     ]),
]):
    lx = Inches(0.5) + col * Inches(6.45)
    add_rect(s, lx, Inches(2.75), Inches(6.1), Inches(3.5), fill_color=GRAY1)
    add_textbox(s, mode,
                lx + Inches(0.2), Inches(2.85), Inches(5.7), Inches(0.38),
                font_size=Pt(11), bold=True, color=WHITE)
    add_rect(s, lx + Inches(0.2), Inches(3.22), Inches(5.7), Inches(0.02),
             fill_color=GRAY3)
    for k, (label, detail) in enumerate(steps_list):
        ty = Inches(3.3) + k * Inches(0.55)
        add_textbox(s, label,
                    lx + Inches(0.25), ty, Inches(1.3), Inches(0.5),
                    font_size=Pt(9), bold=True, color=GRAY4)
        add_textbox(s, detail,
                    lx + Inches(1.55), ty, Inches(4.35), Inches(0.5),
                    font_size=Pt(9), color=GRAY4, font_name=MONO)

# Bottom: signer validations
add_rect(s, Inches(0.5), Inches(6.4), Inches(12.4), Inches(0.62), fill_color=GRAY2)
add_textbox(s,
            "Signer validations on every elevation request:  "
            "(1) allow_sudo must be true for the host  ·  "
            "(2) sudo_user must match allowed_sudo_users regex  ·  "
            "(3) command is shell-quoted to prevent injection",
            Inches(0.65), Inches(6.45), Inches(12.1), Inches(0.52),
            font_size=Pt(9.5), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — AI-ACTION FIREWALL
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "AI-action firewall — three layers\nbetween intent and execution.",
            Inches(0.9), Inches(1.4), Inches(10), Inches(1.4),
            font_size=Pt(32), bold=True, color=WHITE)

layers = [
    ("Phase A · v1.5 → v1.14", "Command Policy",
     "Allow/deny/approval rules evaluated by the signer before signing — a denied command never reaches the server. Composable by group (v1.14): a host inherits the union of its groups' policies plus its own. shell_parse parses POSIX sh so each pipeline stage is evaluated independently.",
     "require_approval flag surfaces commands that need human sign-off."),
    ("Phase B · v1.6", "Human Approval Gate",
     "When a command matches require_approval the broker receives HTTP 202. It polls until a human approves via broker-ctl.",
     "The signer will not issue a certificate without the approval token."),
    ("Phase C · v1.7", "Behaviour Guardrails",
     "Control plane tracks per-subject baselines: rate, known hosts, command vocabulary. Anomalies logged or escalated.",
     "Rate limiting per subject. Exceeding limit returns HTTP 429."),
]

for i, (phase, title, body, note) in enumerate(layers):
    left = Inches(0.5) + i * Inches(4.2)
    top  = Inches(3.15)
    add_rect(s, left, top, Inches(3.9), Inches(3.6), fill_color=GRAY1)
    add_textbox(s, phase,
                left + Inches(0.2), top + Inches(0.18),
                Inches(3.5), Inches(0.5),
                font_size=Pt(9.5), bold=True, color=GRAY4)
    add_textbox(s, title,
                left + Inches(0.2), top + Inches(0.72),
                Inches(3.5), Inches(0.55),
                font_size=Pt(14), bold=True, color=WHITE)
    add_rect(s, left + Inches(0.2), top + Inches(1.32),
             Inches(3.5), Inches(0.025), fill_color=GRAY3)
    add_textbox(s, body,
                left + Inches(0.2), top + Inches(1.45),
                Inches(3.5), Inches(1.55),
                font_size=Pt(10), color=GRAY4)
    add_textbox(s, note,
                left + Inches(0.2), top + Inches(2.98),
                Inches(3.5), Inches(0.6),
                font_size=Pt(9), color=GRAY3, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — COMMAND POLICY: SHELL AST PARSING  (v1.9.2)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "03  Security Controls  ·  v1.9.2")
title_text(s, "Command policy — shell AST parsing", y=Inches(1.4), size=Pt(30))

add_textbox(s,
            "shell_parse: true parses the command as POSIX sh before regex evaluation — "
            "preventing compound-command bypasses like  ps aux && kill -9 1.",
            Inches(0.9), Inches(2.12), Inches(11.5), Inches(0.4),
            font_size=Pt(11), color=GRAY3, italic=True)

# Left column: problem vs solution
left_cards = [
    (GRAY3, "WITHOUT shell_parse",
     'Allowlist ["^ps"] matches "ps aux && kill -9 1000"\nbecause the string starts with "ps".\nThe && is invisible to the regex evaluator.\nsshd runs both commands.'),
    (BLACK, "WITH shell_parse: true",
     'extractCommands("ps aux && kill -9 1000")\n  → ["ps aux", "kill -9 1000"]\ndecideOne("ps aux")       → allow ✓\ndecideOne("kill -9 1000") → no-match ✗\n→ DENIED — cert never issued'),
]
for i, (hc, htxt, body) in enumerate(left_cards):
    cy = Inches(2.7) + i * Inches(2.05)
    add_rect(s, Inches(0.5), cy, Inches(6.0), Inches(1.85),
             fill_color=WHITE, line_color=GRAY4, line_width=Pt(0.75))
    add_label_in_rect(s, htxt, Inches(0.5), cy, Inches(6.0), Inches(0.45),
                      fill_color=hc, text_color=WHITE, font_size=Pt(11), bold=True)
    add_textbox(s, body, Inches(0.7), cy + Inches(0.55),
                Inches(5.7), Inches(1.25),
                font_size=Pt(9.5), color=GRAY2, font_name=MONO)

# Right column: rejected nodes + config
rx = Inches(7.0)
add_textbox(s, "Always rejected when shell_parse: true",
            rx, Inches(2.7), Inches(5.8), Inches(0.35),
            font_size=Pt(11), bold=True, color=BLACK)

rejected = [
    ("CmdSubst",  "$(cat /etc/passwd)",    "Arbitrary subshell"),
    ("ProcSubst", "<(cmd)",                "Process substitution"),
    ("ArithmCmd", "$((expr))",             "Arithmetic with side effects"),
    ("Redirect",  "cmd > /etc/cron.d/x",  "File write"),
]
rw = [Inches(1.5), Inches(2.2), Inches(2.0)]
rx2 = [rx, rx + Inches(1.5), rx + Inches(3.7)]
add_label_in_rect(s, "Node",    rx2[0], Inches(3.12), rw[0], Inches(0.36),
                  fill_color=BLACK, text_color=WHITE, font_size=Pt(9), bold=True)
add_label_in_rect(s, "Example", rx2[1], Inches(3.12), rw[1], Inches(0.36),
                  fill_color=BLACK, text_color=WHITE, font_size=Pt(9), bold=True)
add_label_in_rect(s, "Reason",  rx2[2], Inches(3.12), rw[2], Inches(0.36),
                  fill_color=BLACK, text_color=WHITE, font_size=Pt(9), bold=True)

for i, (node, example, reason) in enumerate(rejected):
    bg = WHITE if i % 2 == 0 else GRAY5
    ry = Inches(3.48) + i * Inches(0.42)
    for cell, w, cx in zip([node, example, reason], rw, rx2):
        add_rect(s, cx, ry, w, Inches(0.42), fill_color=bg,
                 line_color=GRAY4, line_width=Pt(0.5))
        add_textbox(s, cell, cx + Inches(0.1), ry + Inches(0.07),
                    w - Inches(0.15), Inches(0.32),
                    font_size=Pt(9), color=GRAY2, font_name=MONO if i > 0 and i < 3 else SANS)

add_textbox(s, "Pipes (|) and sequences (&&, ;) are structurally allowed —\nbut every stage must pass the policy independently.\nNewlines (\\n, \\r) are rejected outright (v1.11.2).",
            rx, Inches(5.28), Inches(5.8), Inches(0.6),
            font_size=Pt(9.5), color=GRAY2, italic=True)

mono_block(s,
           '"command_policy": {\n  "mode": "allowlist",\n  "shell_parse": true,\n'
           '  "allow": [\n    "^ps aux$",\n    "^grep [a-zA-Z0-9_. -]+"\n  ]\n}',
           rx, Inches(5.9), Inches(5.8), Inches(1.32))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — COMPOSABLE COMMAND POLICIES BY GROUP  (v1.14.0)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS  ·  v1.14.0",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Composable command policies\nby group",
            Inches(0.9), Inches(1.4), Inches(11), Inches(1.4),
            font_size=Pt(31), bold=True, color=WHITE)
add_textbox(s,
            "The firewall is no longer per-host only. A named policy library attaches to groups; "
            "a host inherits the composition of every group it belongs to, plus its own inline policy.",
            Inches(0.9), Inches(2.62), Inches(11.6), Inches(0.4),
            font_size=Pt(11.5), color=GRAY4, italic=True)

# Composition pipeline: library → groups → effective policy
model = [
    ("command_policies", "Named library",
     "Reusable allow / deny / require_approval policy definitions, each with a name."),
    ("group_command_policies", "Group → policies",
     "Maps each group to a list of policy names. The reserved group _default applies to every host."),
    ("Effective policy", "Composed & compiled",
     "host = inline command_policy  ∪  all its groups' policies. Resolved and validated at config load."),
]
for i, (tag, title, body) in enumerate(model):
    left = Inches(0.5) + i * Inches(4.2)
    top  = Inches(3.2)
    add_rect(s, left, top, Inches(3.9), Inches(1.75), fill_color=GRAY1)
    add_textbox(s, tag, left + Inches(0.2), top + Inches(0.16),
                Inches(3.5), Inches(0.34), font_size=Pt(10.5), bold=True,
                color=WHITE, font_name=MONO)
    add_textbox(s, title, left + Inches(0.2), top + Inches(0.54),
                Inches(3.5), Inches(0.3), font_size=Pt(10), bold=True, color=GRAY4)
    add_rect(s, left + Inches(0.2), top + Inches(0.9), Inches(3.5), Inches(0.02),
             fill_color=GRAY3)
    add_textbox(s, body, left + Inches(0.2), top + Inches(1.0),
                Inches(3.5), Inches(0.7), font_size=Pt(9.5), color=GRAY4)
    if i < 2:
        add_textbox(s, "→", left + Inches(3.92), top + Inches(0.6),
                    Inches(0.3), Inches(0.5), font_size=Pt(20), bold=True, color=GRAY3)

# Composition rules strip
add_textbox(s, "Composition is additive — strictest wins:",
            Inches(0.5), Inches(5.18), Inches(7), Inches(0.32),
            font_size=Pt(11), bold=True, color=WHITE)
rules = [
    ("deny", "wins — any denylist match blocks"),
    ("allow", "union of every allowlist"),
    ("require_approval", "union — any match escalates"),
    ("shell_parse", "OR — any policy enables it"),
]
for i, (k, v) in enumerate(rules):
    rx = Inches(0.5) + i * Inches(3.1)
    add_label_in_rect(s, k, rx, Inches(5.55), Inches(3.0), Inches(0.34),
                      fill_color=WHITE, text_color=BLACK, font_size=Pt(10), bold=True)
    add_textbox(s, v, rx + Inches(0.05), Inches(5.95), Inches(2.95), Inches(0.5),
                font_size=Pt(9), color=GRAY4)

# Offline inspection
mono_block(s,
           "$ broker-ctl policy explain --host web01 --command 'systemctl restart nginx'\n"
           "  → composed: _default + prod-web + inline      decision: ALLOWED, requires approval",
           Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.55), bg=GRAY1, fg=GRAY5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14b — DYNAMIC POLICY OPERATIONS  (v1.17.0)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS  ·  v1.17.0",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Dynamic policy operations",
            Inches(0.9), Inches(1.4), Inches(11), Inches(1.0),
            font_size=Pt(31), bold=True, color=WHITE)
add_textbox(s,
            "Manage the firewall without abandoning the file as the source of truth: recommend "
            "changes from the audit, apply them with a validated API, and pick up edits automatically.",
            Inches(0.9), Inches(2.34), Inches(11.6), Inches(0.5),
            font_size=Pt(11.5), color=GRAY4, italic=True)

ops = [
    ("policy recommend", "Mine the audit",
     "Read-only advice from what actually ran: promote (run/approved despite a deny), "
     "dead-rule (never matched), friction (repeatedly blocked). It only suggests — a human decides."),
    ("POST /v1/policy", "Validated mutation",
     "broker-ctl policy add/remove over mTLS (auth: reload_callers). Validated before persist, "
     "written atomically, applied in-memory, recorded in the signed audit log. Never doable by the broker."),
    ("auto_reload_seconds", "Watch & reload",
     "Opt-in: the signer polls its config and hot-reloads on change via the same validated, atomic "
     "path — so a GitOps commit or a hand edit applies on its own."),
]
for i, (tag, title, body) in enumerate(ops):
    left = Inches(0.5) + i * Inches(4.2)
    top = Inches(3.25)
    add_rect(s, left, top, Inches(3.9), Inches(2.15), fill_color=GRAY1)
    add_textbox(s, tag, left + Inches(0.2), top + Inches(0.16),
                Inches(3.5), Inches(0.34), font_size=Pt(10.5), bold=True,
                color=WHITE, font_name=MONO)
    add_textbox(s, title, left + Inches(0.2), top + Inches(0.54),
                Inches(3.5), Inches(0.3), font_size=Pt(10), bold=True, color=GRAY4)
    add_rect(s, left + Inches(0.2), top + Inches(0.9), Inches(3.5), Inches(0.02),
             fill_color=GRAY3)
    add_textbox(s, body, left + Inches(0.2), top + Inches(1.0),
                Inches(3.5), Inches(1.05), font_size=Pt(9.5), color=GRAY4)

mono_block(s,
           "$ broker-ctl policy recommend --audit signer_audit.log   ->  PROMOTE / DEAD-RULE / FRICTION\n"
           "$ broker-ctl policy add --host web01 --allow '^systemctl status [a-z0-9_.-]+$'   ->  validated, atomic, audited",
           Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.75), bg=GRAY1, fg=GRAY5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14c — RUNTIME GRANTS & APPROVE-AND-LEARN  (v1.18.0)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS  ·  v1.18.0",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Runtime grants & approve-and-learn",
            Inches(0.9), Inches(1.4), Inches(11), Inches(1.0),
            font_size=Pt(31), bold=True, color=WHITE)
add_textbox(s,
            "A dynamic overlay composed on top of the file baseline at decision time — temporary, "
            "self-expiring, and widen-only: it can loosen, never tighten or invert.",
            Inches(0.9), Inches(2.34), Inches(11.6), Inches(0.5),
            font_size=Pt(11.5), color=GRAY4, italic=True)

grants = [
    ("policy grant", "Temporary widening",
     "Add allow patterns to an allowlist host for a TTL — expires on its own, no file edit. Refused on a "
     "default-allow host (it would invert it). In-memory, audited, operator-only."),
    ("approval allow --learn", "Approve-and-learn",
     "Approving a require_approval command with --learn mints a TTL'd approval WAIVER: the same command runs "
     "without re-approval until expiry. Bound to the exact command + elevation that was approved."),
    ("widen-only", "Enforced, not assumed",
     "Overlays never override a deny, never invert a host, and are honoured only from the right trust tier. "
     "Dropped on restart (fail-safe) — the gate always returns; every mint is in the signed audit log."),
]
for i, (tag, title, body) in enumerate(grants):
    left = Inches(0.5) + i * Inches(4.2)
    top = Inches(3.25)
    add_rect(s, left, top, Inches(3.9), Inches(2.15), fill_color=GRAY1)
    add_textbox(s, tag, left + Inches(0.2), top + Inches(0.16),
                Inches(3.5), Inches(0.34), font_size=Pt(10.5), bold=True,
                color=WHITE, font_name=MONO)
    add_textbox(s, title, left + Inches(0.2), top + Inches(0.54),
                Inches(3.5), Inches(0.3), font_size=Pt(10), bold=True, color=GRAY4)
    add_rect(s, left + Inches(0.2), top + Inches(0.9), Inches(3.5), Inches(0.02),
             fill_color=GRAY3)
    add_textbox(s, body, left + Inches(0.2), top + Inches(1.0),
                Inches(3.5), Inches(1.05), font_size=Pt(9.5), color=GRAY4)

mono_block(s,
           "$ broker-ctl policy grant --host web01 --allow '^systemctl restart nginx$' --ttl 2h\n"
           "$ broker-ctl approval allow <id> --learn --ttl 2h   ->  skip re-approval for 2h (TTL'd waiver)",
           Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.75), bg=GRAY1, fg=GRAY5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — BEHAVIOUR GUARDRAILS DETAIL
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "03  Security Controls")
title_text(s, "Behaviour guardrails — Phase C", y=Inches(1.4), size=Pt(30))

add_textbox(s,
            "The control plane tracks per-subject baselines and flags deviations — without ML, without configuration.",
            Inches(0.9), Inches(2.12), Inches(11.5), Inches(0.4),
            font_size=Pt(12), color=GRAY3, italic=True)

# Three detection signals
signals = [
    ("Rate spike",
     "Subject exceeds rate_limit_per_min requests/minute.",
     "Immediate HTTP 429 (enforce) or audit anomaly (observe)."),
    ("New host",
     "Command targets a host the subject has never accessed.",
     "Escalated to approval gate (enforce) or audited (observe)."),
    ("Novel command",
     "Command's first token has never appeared in the subject's history.",
     "Same escalation path as new host — anomaly logged with fingerprint."),
]

for i, (title, trigger, action) in enumerate(signals):
    lx = Inches(0.5) + i * Inches(4.2)
    add_rect(s, lx, Inches(2.7), Inches(3.9), Inches(2.9), fill_color=WHITE,
             line_color=GRAY4, line_width=Pt(0.75))
    add_label_in_rect(s, title, lx, Inches(2.7), Inches(3.9), Inches(0.48),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(11), bold=True)
    add_textbox(s, "Trigger:", lx + Inches(0.2), Inches(3.32),
                Inches(0.9), Inches(0.3), font_size=Pt(9), bold=True, color=GRAY3)
    add_textbox(s, trigger, lx + Inches(0.2), Inches(3.6),
                Inches(3.5), Inches(0.8), font_size=Pt(10), color=GRAY2)
    add_rect(s, lx + Inches(0.2), Inches(4.35), Inches(3.5), Inches(0.02),
             fill_color=GRAY4)
    add_textbox(s, "Response:", lx + Inches(0.2), Inches(4.42),
                Inches(1.1), Inches(0.3), font_size=Pt(9), bold=True, color=GRAY3)
    add_textbox(s, action, lx + Inches(0.2), Inches(4.7),
                Inches(3.5), Inches(0.75), font_size=Pt(10), color=BLACK, bold=False)

# Modes table
mode_hdrs = ["Mode", "Rate exceeded", "Anomaly detected", "Audit field"]
mode_rows = [
    ["off",      "—",               "—",                      "—"],
    ["observe",  "—",               "anomaly logged",          "anomaly: true"],
    ["enforce",  "HTTP 429",        "escalated to approval",  "outcome: rate-limited / approval-required"],
]
mw = [Inches(1.4), Inches(2.3), Inches(3.0), Inches(5.3)]
mx = [Inches(0.5)]
for w in mw[:-1]:
    mx.append(mx[-1] + w)
MTOP  = Inches(5.78)
MROH  = Inches(0.48)

for j, (h, w, x) in enumerate(zip(mode_hdrs, mw, mx)):
    add_label_in_rect(s, h, x, MTOP, w, Inches(0.38),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(9.5), bold=True)
for i, row in enumerate(mode_rows):
    bg = WHITE if i % 2 == 0 else GRAY5
    y  = MTOP + Inches(0.38) + i * MROH
    for j, (cell, w, x) in enumerate(zip(row, mw, mx)):
        add_rect(s, x, y, w, MROH, fill_color=bg, line_color=GRAY4, line_width=Pt(0.5))
        add_textbox(s, cell, x + Inches(0.1), y + Inches(0.07),
                    w - Inches(0.15), MROH - Inches(0.1),
                    font_size=Pt(9.5),
                    color=BLACK if j == 0 else GRAY2,
                    bold=(j == 0),
                    font_name=MONO if j == 3 else SANS)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — APPROVAL FLOW SEQUENCE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "03  Security Controls")
title_text(s, "Human approval gate — sequence", y=Inches(1.45), size=Pt(30))

# Four actors
actors = ["Broker", "Control Plane", "Signer", "Human (ops)"]
actor_x = [Inches(0.7), Inches(3.7), Inches(6.7), Inches(9.9)]
actor_w = Inches(2.2)
ACTOR_Y = Inches(2.45)

for x, label in zip(actor_x, actors):
    add_label_in_rect(s, label, x, ACTOR_Y, actor_w, Inches(0.55),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10))
    # dashed lifeline
    cx = x + actor_w / 2
    for seg in range(10):
        sy = ACTOR_Y + Inches(0.55) + seg * Inches(0.42)
        add_rect(s, cx - Inches(0.01), sy, Inches(0.02), Inches(0.22),
                 fill_color=GRAY4)

# Sequence steps: (from_idx, to_idx, y, label, response=False)
seq = [
    (0, 1, Inches(3.15), "POST /v1/sign  { command, pubkey, ... }", False),
    (1, 2, Inches(3.58), "POST /v1/sign  { on_behalf_of=broker-1 }", False),
    (2, 1, Inches(4.01), "200  { require_approval=true, cert=nil }", True),
    (1, 0, Inches(4.44), "202  { approval_id='abc-123' }", True),
    (0, 1, Inches(4.87), "GET /v1/sign/result/abc-123  (polling)", False),
    (1, 3, Inches(5.30), "Notify: approval required (log / webhook)", False),
    (3, 1, Inches(5.73), "POST /v1/approvals/abc-123  { allow }", False),
    (1, 2, Inches(6.16), "POST /v1/sign  { approved=true }", False),
    (2, 0, Inches(6.59), "200  { certificate, serial }  →  SSH proceeds", True),
]

for from_i, to_i, y, label, is_resp in seq:
    x1 = actor_x[from_i] + actor_w / 2
    x2 = actor_x[to_i]   + actor_w / 2
    arrow_color = GRAY3 if not is_resp else GRAY4
    add_arrow(s, x1, y, x2, y, color=arrow_color, width=Pt(1.5 if not is_resp else 1.0))
    mid = (min(x1, x2) + max(x1, x2)) / 2
    add_textbox(s, label,
                min(x1, x2) + Inches(0.1), y - Inches(0.32),
                abs(x2 - x1) - Inches(0.15), Inches(0.3),
                font_size=Pt(8),
                color=GRAY2 if not is_resp else GRAY3,
                bold=(not is_resp),
                align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — TEAMS NOTIFICATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS  ·  v1.8.0",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Approval requests land\nin Microsoft Teams.",
            Inches(0.9), Inches(1.38), Inches(11.5), Inches(1.1),
            font_size=Pt(30), bold=True, color=WHITE)

# ── Mock Adaptive Card ──
card_x, card_y = Inches(0.7), Inches(2.95)
card_w, card_h = Inches(5.5), Inches(3.5)
add_rect(s, card_x, card_y, card_w, card_h, fill_color=WHITE)
add_rect(s, card_x, card_y, card_w, Inches(0.55), fill_color=GRAY2)
add_textbox(s, "SSH Broker · Approval required",
            card_x + Inches(0.2), card_y + Inches(0.13),
            card_w - Inches(0.4), Inches(0.3),
            font_size=Pt(11), bold=True, color=WHITE)

card_facts = [
    ("Requested by", "alice@contoso.com"),
    ("Host",         "db01.prod:22"),
    ("Command",      "pg_dump mydb"),
    ("Elevation",    "sudo -n"),
    ("Approval ID",  "abc-123"),
]
for k, (label, val) in enumerate(card_facts):
    fy = card_y + Inches(0.75) + k * Inches(0.42)
    add_textbox(s, label, card_x + Inches(0.25), fy,
                Inches(1.7), Inches(0.35), font_size=Pt(10), bold=True, color=GRAY2)
    add_textbox(s, val, card_x + Inches(2.0), fy,
                Inches(3.3), Inches(0.35), font_size=Pt(10), color=BLACK)

btn_y = card_y + card_h - Inches(0.65)
add_label_in_rect(s, "View request  ↗",
                  card_x + Inches(0.25), btn_y, Inches(2.2), Inches(0.42),
                  fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)

# ── Right column: how it works ──
rx = Inches(6.7)
teams_points = [
    ("Enable with one line",
     'notifier: "teams" in control-plane.json — reuses the existing webhook_url (Power Automate / M365 connector).'),
    ("Two card formats",
     "Adaptive Card v1.4 (workflow, recommended) or MessageCard (legacy tenants)."),
    ("Deep link to the request",
     "approval_url_template embeds a 'View request' button pointing at {id}."),
    ("Today: one-way notify",
     "Approval still runs via broker-ctl approval allow <id>. Action.Submit buttons are not supported over a simple Incoming Webhook — see next slide."),
]
for k, (h, b) in enumerate(teams_points):
    py = Inches(3.0) + k * Inches(0.92)
    add_rect(s, rx, py + Inches(0.06), Inches(0.08), Inches(0.08), fill_color=WHITE)
    add_textbox(s, h, rx + Inches(0.22), py, Inches(6.0), Inches(0.32),
                font_size=Pt(12), bold=True, color=WHITE)
    add_textbox(s, b, rx + Inches(0.22), py + Inches(0.33), Inches(6.0), Inches(0.55),
                font_size=Pt(9.5), color=GRAY4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — EXTENSIBLE NOTIFIERS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "03  Security Controls")
title_text(s, "Notification architecture —\nextensible by design", y=Inches(1.38), size=Pt(28))

add_textbox(s,
            "Any component that implements Notifier.Notify(Approval) can receive alerts. "
            "Today: three built-in notifiers. Future: multi-channel arrays.",
            Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.45),
            font_size=Pt(11), color=GRAY3, italic=True)

# Three current notifiers
notifiers = [
    ("LogNotifier",     "log",      "Writes to the control-plane audit log.\nAlways active. Zero configuration."),
    ("WebhookNotifier", "webhook",  "HTTP POST JSON payload to any URL.\nWorks with Slack, PagerDuty, custom hooks."),
    ("TeamsNotifier",   "teams",    "Adaptive Card v1.4 (workflow) or MessageCard.\nIncoming Webhook — Power Automate / M365."),
]
for i, (cls, key, desc) in enumerate(notifiers):
    lx = Inches(0.5) + i * Inches(4.2)
    add_rect(s, lx, Inches(2.85), Inches(3.9), Inches(2.5),
             fill_color=WHITE, line_color=GRAY4, line_width=Pt(0.75))
    add_label_in_rect(s, cls, lx, Inches(2.85), Inches(3.9), Inches(0.5),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(11), bold=True)
    add_textbox(s, f'notifier: "{key}"',
                lx + Inches(0.2), Inches(3.5),
                Inches(3.5), Inches(0.32),
                font_size=Pt(10), color=GRAY3, italic=False, font_name=MONO)
    add_textbox(s, desc,
                lx + Inches(0.2), Inches(3.9),
                Inches(3.5), Inches(1.3),
                font_size=Pt(10), color=GRAY2)

# Current config vs future config
add_rect(s, Inches(0.5), Inches(5.55), Inches(12.35), Inches(0.02), fill_color=GRAY4)

col1x, col2x = Inches(0.5), Inches(6.7)
add_textbox(s, "Current config (single notifier)",
            col1x, Inches(5.68), Inches(5.8), Inches(0.3),
            font_size=Pt(10), bold=True, color=GRAY2)
add_textbox(s, "Future config (multi-channel — Phase 2 design)",
            col2x, Inches(5.68), Inches(6.1), Inches(0.3),
            font_size=Pt(10), bold=True, color=GRAY2)

mono_block(s,
           '"approval": {\n  "notifier": "teams",\n  "webhook_url": "...",\n  "teams_format": "workflow"\n}',
           col1x, Inches(6.02), Inches(5.8), Inches(1.2))
mono_block(s,
           '"notifiers": [\n  { "type": "teams", "webhook_url": "..." },\n  { "type": "log" }\n],\n"approval_channels": [\n  { "type": "mtls-cli" },\n  { "type": "teams-bot", "entra_group": "ssh-approvers" }\n]',
           col2x, Inches(6.02), Inches(6.1), Inches(1.2))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — TEAMS APPROVAL BRIDGE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "03  Security Controls")
title_text(s, "Roadmap — approve directly from Teams", y=Inches(1.4), size=Pt(28))

add_textbox(s, "Phase 2: a bridge turns the one-way notification into a closed approval loop with real identity.",
            Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.4),
            font_size=Pt(12), color=GRAY3, italic=True)

# ── Flow diagram ──
FY = Inches(2.65)
FBW, FBH = Inches(3.5), Inches(0.9)
bridge_flow = [
    (Inches(0.5),  "Microsoft Teams\n(Approve / Deny button)",      GRAY2, WHITE),
    (Inches(4.9),  "Approval Bridge\n(validates Entra JWT + group)", GRAY3, WHITE),
    (Inches(9.3),  "Control Plane\nPOST /v1/approvals/{id}",         BLACK, WHITE),
]
for left, label, fc, tc in bridge_flow:
    add_label_in_rect(s, label, left, FY, FBW, FBH,
                      fill_color=fc, text_color=tc, font_size=Pt(10))
bridge_lbls = ["HttpPOST + Entra token", "mTLS (approver cert)"]
for i in range(len(bridge_flow) - 1):
    x1 = bridge_flow[i][0] + FBW
    x2 = bridge_flow[i + 1][0]
    my = FY + FBH / 2
    add_arrow(s, x1, my, x2, my, color=GRAY3, width=Pt(2))
    add_textbox(s, bridge_lbls[i], x1 + Inches(0.02), FY - Inches(0.3),
                x2 - x1 - Inches(0.02), Inches(0.28),
                font_size=Pt(8), color=GRAY3, align=PP_ALIGN.CENTER)

add_textbox(s, "Attribution closes the loop: DecidedBy = the approver's real Entra UPN.",
            Inches(0.5), FY + FBH + Inches(0.12), Inches(12.3), Inches(0.32),
            font_size=Pt(10), color=GRAY2, italic=True, align=PP_ALIGN.CENTER)

# ── Trade-off table ──
bridge_hdrs = ["Option", "Approach", "Attribution", "Cost"]
bridge_rows = [
    ["A — Bot + bridge",
     "Bot Framework + cmd/approval-bridge + Entra group check",
     "Real (approver UPN)", "High"],
    ["B — Signed token",
     "Control plane embeds an HMAC in the button URL; Teams POSTs with no cert",
     "None (anonymous)", "Low"],
    ["C — Notify only (today)",
     "Teams notifies; approver runs broker-ctl approval allow",
     "mTLS CN of approver", "Zero (shipped)"],
]
bw = [Inches(2.6), Inches(5.6), Inches(2.5), Inches(1.6)]
bx = [Inches(0.5)]
for w in bw[:-1]:
    bx.append(bx[-1] + w)
BTT  = Inches(4.3)
BTRH = Inches(0.6)
for j, (h, w, x) in enumerate(zip(bridge_hdrs, bw, bx)):
    add_label_in_rect(s, h, x, BTT, w, Inches(0.4),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)
for i, row in enumerate(bridge_rows):
    bg = WHITE if i % 2 == 0 else GRAY5
    y  = BTT + Inches(0.4) + i * BTRH
    for j, (cell, w, x) in enumerate(zip(row, bw, bx)):
        add_rect(s, x, y, w, BTRH, fill_color=bg, line_color=GRAY4, line_width=Pt(0.5))
        add_textbox(s, cell, x + Inches(0.1), y + Inches(0.08),
                    w - Inches(0.15), BTRH - Inches(0.12),
                    font_size=Pt(9.5), color=BLACK if j == 0 else GRAY2,
                    bold=(j == 0),
                    align=PP_ALIGN.CENTER if j == 3 else PP_ALIGN.LEFT)

add_textbox(s,
            "Recommended: Option A for real-identity approval from Teams / mobile; Option C (today) is sufficient for many deployments.",
            Inches(0.5), Inches(6.62), Inches(12.4), Inches(0.4),
            font_size=Pt(9.5), color=GRAY3, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — RBAC
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "03  Security Controls")
title_text(s, "Multi-layer RBAC", y=Inches(1.5), size=Pt(32))

headers = ["Layer", "Enforcement point", "What it controls"]
rows = [
    ["mTLS client cert (CN)", "Signer", "Which host groups a broker instance can access"],
    ["Host groups", "Signer — GET /v1/hosts", "Broker only receives hosts in its allowed groups"],
    ["OIDC user groups", "Signer — POST /v1/sign", "End-user identity from token propagated to RBAC"],
    ["allow_sudo / allowed_sudo_users", "Signer — policy", "Whether and as whom elevation is permitted"],
    ["command_policy", "Signer — policy", "Allow / deny / require_approval per host"],
    ["Behaviour baseline", "Control plane", "Anomaly detection per subject (rate, host, command)"],
]

col_widths = [Inches(2.8), Inches(2.5), Inches(5.8)]
col_starts = [Inches(0.5), Inches(3.3), Inches(5.8)]
TABLE_TOP = Inches(2.85)
ROW_H     = Inches(0.5)

for j, (h, w) in enumerate(zip(headers, col_widths)):
    add_label_in_rect(s, h, col_starts[j], TABLE_TOP, w, Inches(0.42),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)

for i, row in enumerate(rows):
    bg = WHITE if i % 2 == 0 else GRAY5
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        add_rect(s, col_starts[j],
                 TABLE_TOP + Inches(0.42) + i * ROW_H, w, ROW_H,
                 fill_color=bg, line_color=GRAY4, line_width=Pt(0.5))
        add_textbox(s, cell,
                    col_starts[j] + Inches(0.1),
                    TABLE_TOP + Inches(0.42) + i * ROW_H + Inches(0.08),
                    w - Inches(0.15), ROW_H - Inches(0.1),
                    font_size=Pt(10), color=GRAY2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — AUDIT TRAIL
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Three correlated logs.\nCryptographically chained.",
            Inches(0.9), Inches(1.4), Inches(9), Inches(1.4),
            font_size=Pt(34), bold=True, color=WHITE)

logs = [
    ("Signer log\n(signer_audit.log)",
     ["Certificate issued / denied",
      "caller · host:port · user · principal",
      "elevation · PTY · serial"]),
    ("Broker log\n(audit.log)",
     ["Command executed / denied",
      "caller · host · command · exit_code",
      "serial · session_id · elevation · PTY"]),
    ("sshd log\n(/var/log/auth.log)",
     ["Accepted certificate ID",
      "\"agent=... host=... elev=sudo:root pty=1\"",
      "serial XXXX"]),
]

for i, (title, items) in enumerate(logs):
    left = Inches(0.5) + i * Inches(4.2)
    add_rect(s, left, Inches(3.1), Inches(3.9), Inches(3.5), fill_color=GRAY1)
    add_textbox(s, title, left + Inches(0.2), Inches(3.28),
                Inches(3.5), Inches(0.7), font_size=Pt(12), bold=True, color=WHITE)
    add_rect(s, left + Inches(0.2), Inches(3.95),
             Inches(3.5), Inches(0.025), fill_color=GRAY3)
    for k, item in enumerate(items):
        add_textbox(s, "– " + item,
                    left + Inches(0.25), Inches(4.05) + k * Inches(0.65),
                    Inches(3.4), Inches(0.6), font_size=Pt(10), color=GRAY4)

add_rect(s, Inches(0.5), Inches(6.45), Inches(12.35), Inches(0.38), fill_color=GRAY2)
add_textbox(s,
            "All three logs share a common serial — enabling full end-to-end correlation of every operation.",
            Inches(0.7), Inches(6.47), Inches(12.0), Inches(0.34),
            font_size=Pt(10), color=GRAY4, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — AUDIT TRAIL LIVE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Audit trail — what the log looks like",
            Inches(0.9), Inches(1.38), Inches(11), Inches(0.7),
            font_size=Pt(28), bold=True, color=WHITE)
add_textbox(s, "broker-ctl audit show --log audit.log --json | jq .",
            Inches(0.9), Inches(2.1), Inches(11), Inches(0.38),
            font_size=Pt(10), color=GRAY3, italic=True, font_name=MONO)

entries = [
    ('{ "serial":1042, "ts":"2026-06-08T14:02:11Z", "outcome":"issued",\n'
     '  "caller":"mcp-stdio", "host":"web01.prod:22", "user":"deploy",\n'
     '  "command":"systemctl status nginx", "exit_code":0,\n'
     '  "elevation":"", "pty":false, "policy_rule":"",\n'
     '  "session_id":"", "prev_hash":"a3f1..." }',
     "ISSUED — one-shot, no elevation"),
    ('{ "serial":1043, "ts":"2026-06-08T14:03:45Z", "outcome":"issued",\n'
     '  "caller":"alice@contoso.com", "host":"db01.prod:22", "user":"deploy",\n'
     '  "command":"pg_dump mydb", "exit_code":0,\n'
     '  "elevation":"sudo -n", "pty":false, "policy_rule":"allow:pg_dump.*",\n'
     '  "session_id":"sess-7f2a", "prev_hash":"b8c2..." }',
     "ISSUED — sudo elevation, command matched allow rule"),
    ('{ "serial":1044, "ts":"2026-06-08T14:04:02Z", "outcome":"denied",\n'
     '  "caller":"alice@contoso.com", "host":"db01.prod:22", "user":"deploy",\n'
     '  "command":"rm -rf /var/lib/postgres", "exit_code":-1,\n'
     '  "elevation":"", "pty":false, "policy_rule":"deny:rm.*",\n'
     '  "session_id":"", "prev_hash":"c9d3..." }',
     "DENIED — command matched deny rule, cert never issued"),
]

for i, (entry, caption) in enumerate(entries):
    top = Inches(2.65) + i * Inches(1.52)
    add_rect(s, Inches(0.5), top, Inches(9.0), Inches(1.35), fill_color=GRAY1)
    add_textbox(s, entry,
                Inches(0.65), top + Inches(0.1),
                Inches(8.7), Inches(1.2),
                font_size=Pt(8.5), color=GRAY4, font_name=MONO)
    add_rect(s, Inches(9.6), top, Inches(3.25), Inches(1.35),
             fill_color=GRAY2, line_color=GRAY3, line_width=Pt(0.5))
    add_textbox(s, caption,
                Inches(9.75), top + Inches(0.35),
                Inches(3.0), Inches(0.65),
                font_size=Pt(10), color=WHITE, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — SESSION RECORDING  (v1.10.0)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS  ·  v1.10.0",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Session recording —\na complete picture of what happened.",
            Inches(0.9), Inches(1.38), Inches(10), Inches(1.1),
            font_size=Pt(28), bold=True, color=WHITE)

# Left column — how it works
lx = Inches(0.5)
points = [
    ("shell and pty sessions only",
     "One .cast file per session, named <session_id>.cast. "
     "exec-mode and one-shot ssh_execute are not recorded "
     "(output already in the MCP response and audit log)."),
    ("Three streams captured",
     "\"i\" input — every command the agent typed, before it reached the shell.\n"
     "\"o\" output — stdout, or the merged PTY stream.\n"
     "\"e\" stderr — non-PTY sessions only."),
    ("session_id is the index",
     "Filename matches audit log field session_id. "
     "Use broker-ctl audit show to find recordings by agent, host or time window."),
]
for k, (title, body) in enumerate(points):
    ty = Inches(2.75) + k * Inches(1.42)
    add_rect(s, lx, ty, Inches(0.08), Inches(0.08), fill_color=WHITE)
    add_textbox(s, title, lx + Inches(0.22), ty - Inches(0.05),
                Inches(5.8), Inches(0.32),
                font_size=Pt(12), bold=True, color=WHITE)
    add_textbox(s, body, lx + Inches(0.22), ty + Inches(0.3),
                Inches(5.8), Inches(0.95),
                font_size=Pt(9.5), color=GRAY4)

# Right column — file format + commands
rx = Inches(7.0)
add_textbox(s, "ASCIIcast v2 — self-describing file",
            rx, Inches(2.7), Inches(5.8), Inches(0.35),
            font_size=Pt(11), bold=True, color=WHITE)

mono_block(s,
    '{"version":2,"width":220,"height":40,\n'
    ' "title":"session a3f1 — alice@web01",\n'
    ' "ssh_broker":{"session_id":"a3f1","caller":"alice",\n'
    '               "host":"web01","serial":1042}}\n'
    '[0.000, "i", "df -h /\\n"]\n'
    '[0.012, "o", "Filesystem  Size  Used Avail\\n"]\n'
    '[1.244, "i", "uptime\\n"]\n'
    '[1.251, "o", " 14:00:02 up 3 days\\n"]',
    rx, Inches(3.12), Inches(5.8), Inches(1.75))

add_textbox(s, "Enable & replay",
            rx, Inches(5.0), Inches(5.8), Inches(0.32),
            font_size=Pt(11), bold=True, color=WHITE)

mono_block(s,
    '# config.json\n'
    '"session_recording_dir": "/recordings"\n\n'
    '# play back with full timing\n'
    'asciinema play /recordings/a3f1b2c4.cast\n\n'
    '# extract only what the agent typed\n'
    'jq -r \'select(type=="array" and .[1]=="i") | .[2]\'\\\n'
    '  /recordings/a3f1b2c4.cast',
    rx, Inches(5.4), Inches(5.8), Inches(1.72))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — MULTI-CA AND AZURE KEY VAULT  (v1.11.0)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "03  Security Controls  ·  v1.11.0")
title_text(s, "Multi-CA — one CA key per host group", y=Inches(1.38), size=Pt(28))

add_textbox(s,
            "Each host group can be signed by its own CA key — local PEM or Azure Key Vault HSM. "
            "A compromised CA key is now scoped to its group, not the entire fleet.",
            Inches(0.9), Inches(2.12), Inches(11.5), Inches(0.45),
            font_size=Pt(11), color=GRAY3, italic=True)

# Left column: before / after
left_cards = [
    (GRAY3, "BEFORE v1.11.0 — SINGLE GLOBAL CA",
     "One ca_key for every host.\nA compromised CA = every host is at risk.\nKey must live on disk as a PEM file."),
    (BLACK, "NOW — CA PER HOST GROUP",
     "Each group (prod-web, prod-db…) has its own CA key.\nCA keys can live in Azure Key Vault (HSM-backed).\nA compromised key affects only its group."),
]
for i, (hc, htxt, body) in enumerate(left_cards):
    cy = Inches(2.72) + i * Inches(1.95)
    add_rect(s, Inches(0.5), cy, Inches(6.1), Inches(1.78),
             fill_color=WHITE, line_color=GRAY4, line_width=Pt(0.75))
    add_label_in_rect(s, htxt, Inches(0.5), cy, Inches(6.1), Inches(0.45),
                      fill_color=hc, text_color=WHITE, font_size=Pt(10), bold=True)
    add_textbox(s, body, Inches(0.7), cy + Inches(0.55),
                Inches(5.75), Inches(1.12),
                font_size=Pt(10), color=GRAY2)

# Right column: config + CA selection diagram
rx = Inches(7.0)
add_textbox(s, "signer.json  —  ca_keys",
            rx, Inches(2.7), Inches(5.8), Inches(0.35),
            font_size=Pt(11), bold=True, color=BLACK)

mono_block(s,
           '"ca_keys": {\n'
           '  "_default": { "type": "pem",\n'
           '    "path": "pki/ssh_ca" },\n'
           '  "prod-web": { "type": "akv",\n'
           '    "vault_url": "https://vault.azure.net/",\n'
           '    "key_name": "ssh-ca-web" },\n'
           '  "prod-db":  { "type": "akv",\n'
           '    "key_name": "ssh-ca-db" }\n'
           '}',
           rx, Inches(3.1), Inches(5.8), Inches(1.85))

add_textbox(s, "CA selection — first matching group wins",
            rx, Inches(5.1), Inches(5.8), Inches(0.32),
            font_size=Pt(10), bold=True, color=BLACK)

sel_rows = [
    ("web01  (groups: [prod-web])",  "→  signed by  ssh-ca-web  (AKV)"),
    ("db01   (groups: [prod-db])",   "→  signed by  ssh-ca-db   (AKV)"),
    ("dev01  (groups: [dev])",        "→  signed by  pki/ssh_ca  (PEM default)"),
]
for k, (host, ca) in enumerate(sel_rows):
    ry = Inches(5.48) + k * Inches(0.43)
    bg = WHITE if k % 2 == 0 else GRAY5
    add_rect(s, rx, ry, Inches(5.8), Inches(0.4),
             fill_color=bg, line_color=GRAY4, line_width=Pt(0.5))
    add_textbox(s, host, rx + Inches(0.1), ry + Inches(0.07),
                Inches(2.8), Inches(0.3), font_size=Pt(9), color=GRAY2, font_name=MONO)
    add_textbox(s, ca, rx + Inches(2.9), ry + Inches(0.07),
                Inches(2.8), Inches(0.3), font_size=Pt(9), bold=True, color=BLACK, font_name=MONO)

add_textbox(s, "Backward compatible: ca_key (legacy PEM) still works; ca_keys[\"_default\"] takes precedence when both are set.",
            Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.32),
            font_size=Pt(9), color=GRAY3, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — HARDENING: FAIL-CLOSED BY DEFAULT  (v1.11.2 / v1.12.0)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS  ·  v1.11.2 / v1.12.0",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Hardening — fail-closed by default",
            Inches(0.9), Inches(1.38), Inches(11.5), Inches(0.7),
            font_size=Pt(30), bold=True, color=WHITE)
add_textbox(s, "A misconfiguration should deny access, never silently grant it. Six controls added across the latest releases.",
            Inches(0.9), Inches(2.15), Inches(11.5), Inches(0.4),
            font_size=Pt(11), color=GRAY4, italic=True)

hardening = [
    ("v1.11.2", "Fail-closed OIDC groups",
     "A token missing the configured groups claim is rejected (401), never accepted unrestricted. A claim typo or IdP change cannot silently disable per-user RBAC."),
    ("v1.11.2", "Fail-closed token age",
     "With max_token_age_seconds set, a token without a numeric iat is rejected — its age cannot be established, so it is not trusted."),
    ("v1.11.2", "Newline rejection",
     "\\n / \\r in a one-shot command is rejected by the signer on every host — a newline could smuggle extra lines past the regex firewall (ps\\nrm -rf /)."),
    ("v1.12.0", "Host list scoped to user",
     "ssh_list_servers is filtered by the end user's OIDC groups — the model is only shown hosts it can actually sign for, matching the signer's check."),
    ("v1.12.0", "Bounded approval state",
     "Terminal approval requests are purged 2×TTL after creation — no unbounded memory growth in the control plane registry."),
    ("v1.12.0", "Uniform DoS limits",
     "Request-body caps and HTTP timeouts on every frontend; cmd/broker was brought in line with the signer and HTTP MCP server."),
]

for i, (ver, title, body) in enumerate(hardening):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + col * Inches(4.2)
    top  = Inches(2.75) + row * Inches(1.95)
    add_rect(s, left, top, Inches(3.9), Inches(1.78), fill_color=GRAY1)
    add_textbox(s, ver,
                left + Inches(0.2), top + Inches(0.14),
                Inches(3.5), Inches(0.3),
                font_size=Pt(9), bold=True, color=GRAY4)
    add_textbox(s, title,
                left + Inches(0.2), top + Inches(0.42),
                Inches(3.5), Inches(0.4),
                font_size=Pt(12.5), bold=True, color=WHITE)
    add_rect(s, left + Inches(0.2), top + Inches(0.82),
             Inches(3.5), Inches(0.02), fill_color=GRAY3)
    add_textbox(s, body,
                left + Inches(0.2), top + Inches(0.92),
                Inches(3.5), Inches(0.8),
                font_size=Pt(9), color=GRAY4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — ADVERSARIAL SECURITY REVIEW  (v1.13.0)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "03  Security Controls  ·  v1.13.0")
title_text(s, "Adversarial (red-team) review", y=Inches(1.4), size=Pt(30))

add_textbox(s,
            "A red-team pass over authentication, RBAC, privilege escalation, the command firewall "
            "and audit integrity. Two high-severity bypasses closed, plus several hardening fixes — each with a regression test.",
            Inches(0.9), Inches(2.12), Inches(11.6), Inches(0.5),
            font_size=Pt(11), color=GRAY3, italic=True)

rt_headers = ["Finding", "What an attacker could do", "Sev"]
rt_rows = [
    ["role=bastion firewall bypass",
     "A compromised broker requested role=bastion on a policy host to get an unrestricted cert — no force-command, port-forwarding on. Non-target roles are now rejected on policy hosts.",
     "HIGH"],
    ["Deny-all RBAC collapsed open",
     "Empty OIDC groups (deny-all) were dropped on the wire by omitempty and read as unrestricted — the inverse decision. The wire field now round-trips [] faithfully.",
     "HIGH"],
    ["GET /v1/hosts ignored allowed_callers",
     "A broker excluded from a host still received its addr/user/host_key. The host list now applies the same per-caller authorization as signing.",
     "MED"],
    ["Approval hid sudo elevation",
     "An approver could authorize a benign-looking command unaware it would run as root. broker-ctl and the notifiers now show elevation=sudo:<user>.",
     "MED"],
    ["KeyID control-char injection",
     "A newline in end_user could splice forged lines into the host's sshd auth.log. The signer now rejects control characters in caller / end_user.",
     "MED"],
    ["Audit rotation now verifiable",
     "audit verify --all cross-links rotated segments (each prev_hash chains to the previous file) so a dropped or truncated segment is detectable.",
     "FIX"],
]

rcol_w = [Inches(3.2), Inches(7.5), Inches(0.95)]
rcol_x = [Inches(0.5), Inches(3.7), Inches(11.2)]
RTOP   = Inches(2.78)
RROW_H = Inches(0.66)

for j, (h, w, x) in enumerate(zip(rt_headers, rcol_w, rcol_x)):
    add_label_in_rect(s, h, x, RTOP, w, Inches(0.4),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)

for i, row in enumerate(rt_rows):
    bg = WHITE if i % 2 == 0 else GRAY5
    y  = RTOP + Inches(0.4) + i * RROW_H
    for j, (cell, w, x) in enumerate(zip(row, rcol_w, rcol_x)):
        add_rect(s, x, y, w, RROW_H, fill_color=bg, line_color=GRAY4, line_width=Pt(0.5))
        if j == 2:
            badge = BLACK if cell == "HIGH" else (GRAY3 if cell == "MED" else GRAY4)
            add_rect(s, x + Inches(0.1), y + Inches(0.16), w - Inches(0.2), Inches(0.34),
                     fill_color=badge)
            add_textbox(s, cell, x + Inches(0.05), y + Inches(0.2),
                        w - Inches(0.1), Inches(0.3), font_size=Pt(8.5), bold=True,
                        color=WHITE, align=PP_ALIGN.CENTER)
        else:
            add_textbox(s, cell, x + Inches(0.1), y + Inches(0.07),
                        w - Inches(0.18), RROW_H - Inches(0.12),
                        font_size=Pt(10 if j == 0 else 9),
                        color=BLACK if j == 0 else GRAY2, bold=(j == 0))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — DEPLOYMENT LOCAL MODE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "04  Deployment — Mode 1")
title_text(s, "Local mode  (stdio MCP)", y=Inches(1.4), size=Pt(30))

add_textbox(s, "The AI model runs on the same machine as the broker. No network exposure, no token required.",
            Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5),
            font_size=Pt(12), color=GRAY3, italic=True)

LOCAL_Y = Inches(3.05)
LOCAL_H = Inches(0.85)
LOCAL_W = Inches(2.2)

local_boxes = [
    (Inches(0.4),  "AI Model\n(Claude / IDE)",       GRAY2, WHITE),
    (Inches(3.1),  "mcp-broker\n(local binary)",      GRAY1, WHITE),
    (Inches(5.8),  "Control Plane\n(optional)",       GRAY3, WHITE),
    (Inches(8.5),  "Signer\n(CA custodian)",          BLACK, WHITE),
]
for left, label, fc, tc in local_boxes:
    add_label_in_rect(s, label, left, LOCAL_Y, LOCAL_W, LOCAL_H,
                      fill_color=fc, text_color=tc, font_size=Pt(10))

local_lbls = ["stdio\nMCP tools", "mTLS HTTPS", "mTLS HTTPS"]
for i in range(len(local_boxes) - 1):
    x1 = local_boxes[i][0] + LOCAL_W
    x2 = local_boxes[i + 1][0]
    my = LOCAL_Y + LOCAL_H / 2
    add_arrow(s, x1, my, x2, my, color=GRAY3, width=Pt(2))
    add_textbox(s, local_lbls[i],
                x1 + Inches(0.05), LOCAL_Y - Inches(0.5),
                x2 - x1 - Inches(0.05), Inches(0.45),
                font_size=Pt(8), color=GRAY3, align=PP_ALIGN.CENTER)

add_rect(s, Inches(3.0), Inches(2.65), Inches(7.85), Inches(1.55),
         line_color=GRAY4, line_width=Pt(0.75))
add_textbox(s, "Internal network — mTLS (mutual TLS, client cert CN=broker-1)",
            Inches(3.15), Inches(2.67), Inches(7.5), Inches(0.32),
            font_size=Pt(8), color=GRAY3, italic=True)

add_rect(s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(1.85),
         fill_color=GRAY5, line_color=GRAY4, line_width=Pt(0.75))
add_textbox(s, "SSH Infrastructure  (cert signed by CA, ephemeral — valid 60–120 s)",
            Inches(0.6), Inches(4.75), Inches(9), Inches(0.35),
            font_size=Pt(9), bold=True, color=GRAY3)

mid_bx2 = local_boxes[1][0] + LOCAL_W / 2
add_arrow(s, mid_bx2, LOCAL_Y + LOCAL_H,
             mid_bx2, Inches(4.7), color=GRAY3, width=Pt(1.5))
add_textbox(s, "SSH + ephemeral cert",
            mid_bx2 - Inches(0.9), LOCAL_Y + LOCAL_H + Inches(0.05),
            Inches(1.8), Inches(0.3), font_size=Pt(8), color=GRAY3)

ssh_local = [
    (Inches(0.9),  "Bastion :22"),
    (Inches(3.7),  "Web servers\n(prod-web)"),
    (Inches(6.5),  "DB servers\n(prod-db)"),
    (Inches(9.3),  "Other hosts"),
    (Inches(11.8), "…"),
]
for left, lbl in ssh_local:
    add_label_in_rect(s, lbl, left, Inches(5.15), Inches(2.5), Inches(0.72),
                      fill_color=WHITE, text_color=GRAY2, font_size=Pt(9), bold=False)
    add_rect(s, left, Inches(5.15), Inches(2.5), Inches(0.72),
             line_color=GRAY4, line_width=Pt(0.5))

add_textbox(s, "stdout / stderr / exit_code  →  AI model (via stdio)",
            Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.35),
            font_size=Pt(10), color=GRAY3, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — DEPLOYMENT REMOTE MODE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "04  DEPLOYMENT — MODE 2",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Remote mode  (HTTP MCP + OAuth 2.1 / Entra ID)",
            Inches(0.9), Inches(1.38), Inches(11.5), Inches(0.75),
            font_size=Pt(28), bold=True, color=WHITE)
add_textbox(s, "Multi-user, network-accessible. The AI client authenticates with a bearer token issued by Microsoft Entra ID (Azure AD).",
            Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.45),
            font_size=Pt(11), color=GRAY4, italic=True)

TOP_Y = Inches(2.8)
BOX_H2 = Inches(0.8)
BOX_W2 = Inches(2.1)

add_textbox(s, "CLIENT",
            Inches(0.4), TOP_Y - Inches(0.3), Inches(1.2), Inches(0.28),
            font_size=Pt(8), bold=True, color=GRAY4)
add_rect(s, Inches(0.4), TOP_Y - Inches(0.02), Inches(0.04), BOX_H2 + Inches(0.04),
         fill_color=GRAY3)

client_boxes = [
    (Inches(0.65),  "AI Model\n(Claude / IDE)",  GRAY2, WHITE),
    (Inches(3.3),   "MCP Client\n(HTTP+Bearer)", GRAY1, WHITE),
    (Inches(6.2),   "Microsoft\nEntra ID",       GRAY2, WHITE),
]
for left, label, fc, tc in client_boxes:
    add_label_in_rect(s, label, left, TOP_Y, BOX_W2, BOX_H2,
                      fill_color=fc, text_color=tc, font_size=Pt(10))

add_arrow(s, client_boxes[0][0] + BOX_W2, TOP_Y + BOX_H2 / 2,
             client_boxes[1][0],            TOP_Y + BOX_H2 / 2,
             color=GRAY3, width=Pt(2))
add_textbox(s, "tool call",
            client_boxes[0][0] + BOX_W2 + Inches(0.05), TOP_Y - Inches(0.32),
            Inches(1.2), Inches(0.28), font_size=Pt(8), color=GRAY3, align=PP_ALIGN.CENTER)

add_arrow(s, client_boxes[1][0] + BOX_W2, TOP_Y + BOX_H2 / 2,
             client_boxes[2][0],            TOP_Y + BOX_H2 / 2,
             color=GRAY4, width=Pt(1.5))
add_textbox(s, "Auth Code + PKCE (RFC 6749 / OAuth 2.1)",
            client_boxes[1][0] + BOX_W2 + Inches(0.05), TOP_Y - Inches(0.35),
            Inches(2.65), Inches(0.32), font_size=Pt(8), color=GRAY4, align=PP_ALIGN.CENTER)
add_textbox(s, "JWT  (sub, groups, iat)",
            client_boxes[2][0], TOP_Y + BOX_H2 + Inches(0.07),
            Inches(2.1), Inches(0.32), font_size=Pt(8), color=GRAY4, align=PP_ALIGN.CENTER)

BOT_Y = Inches(4.5)
add_textbox(s, "SERVER",
            Inches(0.4), BOT_Y - Inches(0.3), Inches(1.2), Inches(0.28),
            font_size=Pt(8), bold=True, color=GRAY4)
add_rect(s, Inches(0.4), BOT_Y - Inches(0.02), Inches(0.04), BOX_H2 + Inches(0.04),
         fill_color=GRAY3)

server_boxes = [
    (Inches(0.65),  "mcp-broker-http\n(HTTPS + OIDC)",  GRAY1, WHITE),
    (Inches(3.6),   "Control Plane\n(PEP)",              GRAY3, WHITE),
    (Inches(6.55),  "Signer\n(CA custodian)",            BLACK, WHITE),
    (Inches(9.5),   "SSH hosts\n(via ephemeral cert)",   GRAY2, WHITE),
]
for left, label, fc, tc in server_boxes:
    add_label_in_rect(s, label, left, BOT_Y, BOX_W2, BOX_H2,
                      fill_color=fc, text_color=tc, font_size=Pt(10))

srv_lbls = ["mTLS HTTPS", "mTLS HTTPS", "SSH + cert\n(60–120 s)"]
for i in range(len(server_boxes) - 1):
    x1 = server_boxes[i][0] + BOX_W2
    x2 = server_boxes[i + 1][0]
    my = BOT_Y + BOX_H2 / 2
    add_arrow(s, x1, my, x2, my, color=GRAY3, width=Pt(2))
    add_textbox(s, srv_lbls[i],
                x1 + Inches(0.05), BOT_Y - Inches(0.38),
                x2 - x1 - Inches(0.05), Inches(0.33),
                font_size=Pt(8), color=GRAY3, align=PP_ALIGN.CENTER)

mid_mcpc_x = client_boxes[1][0] + BOX_W2 / 2
mid_httpb_x = server_boxes[0][0] + BOX_W2 / 2
add_arrow(s, mid_mcpc_x, TOP_Y + BOX_H2,
             mid_httpb_x, BOT_Y, color=GRAY4, width=Pt(1.5))
add_textbox(s, "HTTPS  Authorization: Bearer <JWT>",
            Inches(0.65), TOP_Y + BOX_H2 + Inches(0.1),
            Inches(3.5), Inches(0.32), font_size=Pt(8), color=GRAY4, italic=True)

mid_entra_x = client_boxes[2][0] + BOX_W2 / 2
add_arrow(s, mid_entra_x, BOT_Y,
             mid_entra_x, TOP_Y + BOX_H2, color=GRAY3, width=Pt(1))
add_textbox(s, "JWKS validation\n(go-oidc, no round-trip\nper request)",
            mid_entra_x + Inches(0.1), TOP_Y + BOX_H2 + Inches(0.15),
            Inches(2.2), Inches(0.7), font_size=Pt(8), color=GRAY3, italic=True)

add_rect(s, Inches(9.3), Inches(3.05), Inches(3.65), Inches(0.95), fill_color=GRAY1)
add_textbox(s, "groups claim → RBAC",
            Inches(9.45), Inches(3.08), Inches(3.4), Inches(0.28),
            font_size=Pt(9), bold=True, color=WHITE)
add_textbox(s, "Token groups ∩ host groups\n→ per-user host visibility\nand command policy",
            Inches(9.45), Inches(3.35), Inches(3.4), Inches(0.6),
            font_size=Pt(8), color=GRAY4)

add_textbox(s, "stdout / stderr / exit_code  →  MCP client  →  AI model  (over HTTPS)",
            Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.35),
            font_size=Pt(10), color=GRAY3, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — ENTRA ID — EXECUTIVE INTEGRATION
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "04  Deployment")
title_text(s, "Connecting to Microsoft Entra ID", y=Inches(1.45), size=Pt(30))

add_textbox(s,
            "Three steps for your IT team to enable federated identity for SSH Broker.",
            Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.4),
            font_size=Pt(12), color=GRAY3, italic=True)

cards_entra = [
    ("Step 1 — Azure Portal\n(IT team, ~15 min)",
     [
         "Register a new App in Entra ID (Azure AD).",
         "Set Redirect URI to the broker callback URL.",
         "Add a security group for each access level\n(e.g. ssh-prod-web, ssh-prod-db).",
         "Enable the 'groups' claim in the token manifest.",
         "Share the Tenant ID and Client ID with the broker admin.",
     ]),
    ("Step 2 — Broker configuration\n(broker admin, ~5 min)",
     [
         "Set oauth.issuer to your Entra ID tenant URL.",
         "Set oauth.client_id from the App Registration.",
         "Set oauth.groups_claim = \"groups\".",
         "Configure resource_url (the broker's public HTTPS address).",
         "No client secret needed — validation uses JWKS (public keys only).",
     ]),
    ("Step 3 — Access policy\n(broker admin, ~5 min)",
     [
         "In signer.json, add a 'groups' field to each host:",
         "  web01 → groups: [\"ssh-prod-web\"]",
         "  db01  → groups: [\"ssh-prod-db\"]",
         "Users in the Entra group ssh-prod-web can only\naccess web hosts — enforced by the signer.",
         "Add or remove users from Entra groups to\ngrant or revoke access instantly.",
     ]),
]

for i, (title, items) in enumerate(cards_entra):
    lx = Inches(0.45) + i * Inches(4.25)
    add_rect(s, lx, Inches(2.75), Inches(4.0), Inches(4.0),
             fill_color=WHITE, line_color=GRAY4, line_width=Pt(0.75))
    add_textbox(s, title,
                lx + Inches(0.2), Inches(2.85), Inches(3.6), Inches(0.65),
                font_size=Pt(11), bold=True, color=BLACK)
    add_rect(s, lx + Inches(0.2), Inches(3.48), Inches(3.6), Inches(0.025),
             fill_color=GRAY4)
    for k, item in enumerate(items):
        add_textbox(s, "· " + item,
                    lx + Inches(0.2), Inches(3.58) + k * Inches(0.6),
                    Inches(3.65), Inches(0.58),
                    font_size=Pt(9.5), color=GRAY2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — SECURITY PROPERTIES
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "03  SECURITY CONTROLS",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "What you gain.",
            Inches(0.9), Inches(1.4), Inches(9), Inches(0.9),
            font_size=Pt(38), bold=True, color=WHITE)

props = [
    ("No static credentials",
     "An attacker who compromises the AI process or the broker obtains no reusable secret. Ed25519 keys live only in RAM for the duration of a single operation."),
    ("Blast radius containment",
     "Policy enforcement lives in the signer — a separate process / machine. A compromised broker cannot elevate privileges or access hosts outside its assigned groups."),
    ("Command-level enforcement",
     "The signer bakes force-command into the certificate. sshd enforces it regardless of what the broker claims. No bypass is possible from the broker side."),
    ("Full attribution",
     "Audit logs record AI identity, end-user OIDC subject, host, command, exit code, elevation and PTY — all chained with Ed25519 signatures."),
    ("Human-in-the-loop",
     "High-risk commands trigger an approval gate. The broker polls; a human approves via CLI or API. The certificate is not issued until approval is recorded."),
    ("Behaviour anomaly detection",
     "The control plane tracks per-subject baselines. Unusual rate, unknown host or novel command vocabulary can trigger alerts or automatic escalation to approval."),
]

for i, (title, body) in enumerate(props):
    row = i // 2
    col = i % 2
    left = Inches(0.5) + col * Inches(6.2)
    top  = Inches(2.55) + row * Inches(1.55)
    add_rect(s, left, top, Inches(5.9), Inches(1.4), fill_color=GRAY1)
    add_textbox(s, title,
                left + Inches(0.2), top + Inches(0.15),
                Inches(5.5), Inches(0.45),
                font_size=Pt(12), bold=True, color=WHITE)
    add_textbox(s, body,
                left + Inches(0.2), top + Inches(0.6),
                Inches(5.5), Inches(0.75),
                font_size=Pt(10), color=GRAY4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 31 — DAY-TO-DAY OPERATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "05  OPERATIONS & ROADMAP",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)
add_textbox(s, "Day-to-day operations\nwith broker-ctl",
            Inches(0.9), Inches(1.38), Inches(9), Inches(1.0),
            font_size=Pt(30), bold=True, color=WHITE)

ops_blocks = [
    ("Host management",
     "# Add a host (auto-scan SSH host key)\nbroker-ctl host add --name web01 \\\n  --addr 10.0.0.1:22 --user deploy \\\n  --scan --sudo --pty --groups prod-web\n\n# List configured hosts\nbroker-ctl host list\n\n# Remove a host\nbroker-ctl host remove web01"),
    ("Audit — follow & filter",
     "# Stream the live execution log\nbroker-ctl audit tail --log audit.log\n\n# Filter by host and outcome\nbroker-ctl audit show --log audit.log \\\n  --host db01 --outcome denied\n\n# Verify cryptographic chain\nbroker-ctl audit verify --log audit.log \\\n  --key pki/audit.seed"),
    ("Approvals",
     "# List pending approvals\nbroker-ctl approval list\n\n# Approve a pending command\nbroker-ctl approval allow abc-123\n\n# Deny a pending command\nbroker-ctl approval deny abc-123\n\n# (mTLS to control plane — requires\n#  approver cert)"),
    ("Reload config\n(no restart needed)",
     "# Hot-reload signer.json after editing\n# (updates hosts, max_ttl, ca_key)\n\n# Via SIGHUP (local):\nkill -HUP \"$(cat signer.pid)\"\n\n# Via HTTP (remote):\nbroker-ctl reload\n\n# Broker picks up new hosts in ≤ 30 s"),
]

for i, (title, code) in enumerate(ops_blocks):
    col = i % 2
    row = i // 2
    lx = Inches(0.5) + col * Inches(6.3)
    ty = Inches(2.65) + row * Inches(2.18)
    add_rect(s, lx, ty, Inches(5.9), Inches(2.0), fill_color=GRAY1)
    add_textbox(s, title,
                lx + Inches(0.2), ty + Inches(0.12),
                Inches(5.5), Inches(0.35),
                font_size=Pt(11), bold=True, color=WHITE)
    add_rect(s, lx + Inches(0.2), ty + Inches(0.48), Inches(5.5), Inches(0.02),
             fill_color=GRAY3)
    add_textbox(s, code,
                lx + Inches(0.2), ty + Inches(0.55),
                Inches(5.5), Inches(1.38),
                font_size=Pt(8.5), color=GRAY4, font_name=MONO)

add_textbox(s,
            "v1.15: --config is a global flag (before the subcommand) and every binary reports --version.",
            Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.3),
            font_size=Pt(9), color=GRAY3, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 32 — GAPS TOWARD PRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "05  Operations & Roadmap")
title_text(s, "Gaps toward production", y=Inches(1.45), size=Pt(32))

add_textbox(s, "Known items and their impact before a production deployment.",
            Inches(0.9), Inches(2.22), Inches(11.5), Inches(0.38),
            font_size=Pt(12), color=GRAY3, italic=True)

gap_headers = ["Gap", "Description", "Effort", "Blocks prod?"]
gap_rows = [
    ["HSM / KMS for CA key\n(AWS KMS / GCP)",
     "Azure Key Vault supported (v1.11.0). AWS KMS and GCP Cloud HSM still need implementation via the crypto.Signer interface.",
     "S", "No"],
    ["Key Revocation List (KRL)",
     "No mechanism to invalidate an issued certificate before its TTL expires. Add /v1/revoke endpoint and RevokedKeys in sshd_config.",
     "S", "No"],
    ["Multi-instance session store",
     "Session state is in-process. Running multiple broker replicas for HA requires externalising to Redis with TTL.",
     "M", "Yes (if HA)"],
    ["WORM audit log export",
     "Audit logs are local files. For compliance, stream the signed chain to S3, GCS, Loki or a SIEM in real time.",
     "M", "No"],
]

gcol_w = [Inches(2.4), Inches(6.0), Inches(1.2), Inches(1.85)]
gcol_x = [Inches(0.5), Inches(2.9), Inches(8.9), Inches(10.1)]
GTOP   = Inches(2.8)
GROW_H = Inches(0.72)

for j, (h, w, x) in enumerate(zip(gap_headers, gcol_w, gcol_x)):
    add_label_in_rect(s, h, x, GTOP, w, Inches(0.42),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)

for i, row in enumerate(gap_rows):
    bg = WHITE if i % 2 == 0 else GRAY5
    y  = GTOP + Inches(0.42) + i * GROW_H
    for j, (cell, w, x) in enumerate(zip(row, gcol_w, gcol_x)):
        add_rect(s, x, y, w, GROW_H, fill_color=bg,
                 line_color=GRAY4, line_width=Pt(0.5))
        # "Yes" in Blocks col gets dark highlight
        if j == 3 and cell == "Yes":
            add_rect(s, x + Inches(0.1), y + Inches(0.1),
                     w - Inches(0.2), GROW_H - Inches(0.2),
                     fill_color=GRAY2)
            add_textbox(s, cell, x + Inches(0.1), y + Inches(0.18),
                        w - Inches(0.2), GROW_H - Inches(0.3),
                        font_size=Pt(10), bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER)
        else:
            fc = BLACK if j == 0 else GRAY2
            add_textbox(s, cell, x + Inches(0.1), y + Inches(0.08),
                        w - Inches(0.15), GROW_H - Inches(0.12),
                        font_size=Pt(9.5 if j != 1 else 9),
                        color=fc, bold=(j == 0),
                        align=PP_ALIGN.CENTER if j in (2, 3) else PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 33 — SECURITY LIMITS / NON-GOALS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "05  Operations & Roadmap")
title_text(s, "Security limits — what we don't claim", y=Inches(1.45), size=Pt(30))

add_textbox(s,
            "The honest counterpart to the threat model. Documented in THREAT_MODEL.md and SECURITY.md — "
            "by design, not oversight.",
            Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.4),
            font_size=Pt(12), color=GRAY3, italic=True)

lim_headers = ["Limit / non-goal", "Why it exists", "Mitigation today"]
lim_rows = [
    ["Sessions have no\ncommand firewall",
     "force-command is one-shot only; in a session the command is invisible to the signer at signing time.",
     "Hosts with a command_policy reject sessions outright. Sensitive hosts → one-shot. TTL + source-address + sudoers."],
    ["Behaviour = detection,\nnot containment",
     "The guardrail subject is asserted by the broker — a compromised broker can rotate identity to reset baselines.",
     "The authoritative controls are signer-side policy + the approval gate, which a broker cannot bypass."],
    ["No certificate\nrevocation (KRL)",
     "An issued certificate cannot be invalidated before it expires.",
     "Short TTL (60–120 s) bounds the exposure window. KRL endpoint is on the roadmap."],
    ["callers is\ndefault-open",
     "A broker CN absent from the callers table has no group restriction (backward-compatible).",
     "List every broker CN explicitly; pin sensitive hosts with per-host allowed_callers."],
]

lcol_w = [Inches(2.7), Inches(4.5), Inches(5.1)]
lcol_x = [Inches(0.5), Inches(3.2), Inches(7.7)]
LTOP   = Inches(2.78)
LROW_H = Inches(0.92)

for j, (h, w, x) in enumerate(zip(lim_headers, lcol_w, lcol_x)):
    add_label_in_rect(s, h, x, LTOP, w, Inches(0.42),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)

for i, row in enumerate(lim_rows):
    bg = WHITE if i % 2 == 0 else GRAY5
    y  = LTOP + Inches(0.42) + i * LROW_H
    for j, (cell, w, x) in enumerate(zip(row, lcol_w, lcol_x)):
        add_rect(s, x, y, w, LROW_H, fill_color=bg,
                 line_color=GRAY4, line_width=Pt(0.5))
        fc = BLACK if j == 0 else GRAY2
        add_textbox(s, cell, x + Inches(0.12), y + Inches(0.1),
                    w - Inches(0.2), LROW_H - Inches(0.15),
                    font_size=Pt(9), color=fc, bold=(j == 0))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 34 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "05  Operations & Roadmap")
title_text(s, "Competitive landscape", y=Inches(1.5), size=Pt(32))

comp_headers = ["Tool", "Ephemeral certs", "MCP native", "AI command policy", "HSM / KMS"]
comp_rows = [
    ["SSH Broker (this project)", "✓", "✓", "✓", "✓ AKV (v1.11)"],
    ["Teleport",                  "✓", "✓ (2025)",  "—", "✓ (cluster)"],
    ["HashiCorp Vault SSH",       "✓", "—", "—", "✓ (managed keys)"],
    ["Smallstep SSH CA",          "✓", "—", "—", "~"],
    ["StrongDM",                  "✗", "—", "—", "✓"],
    ["ssh-mcp",                   "✗", "✓", "—", "✗"],
]

ccol_w = [Inches(3.2), Inches(2.0), Inches(1.8), Inches(2.4), Inches(1.8)]
ccol_x = [Inches(0.5)]
for w in ccol_w[:-1]:
    ccol_x.append(ccol_x[-1] + w)

CTOP  = Inches(3.0)
CROW_H = Inches(0.5)

for j, (h, w, x) in enumerate(zip(comp_headers, ccol_w, ccol_x)):
    add_label_in_rect(s, h, x, CTOP, w, Inches(0.42),
                      fill_color=BLACK, text_color=WHITE, font_size=Pt(10), bold=True)

for i, row in enumerate(comp_rows):
    bg = GRAY5 if i == 0 else (WHITE if i % 2 == 1 else GRAY5)
    for j, (cell, w, x) in enumerate(zip(row, ccol_w, ccol_x)):
        add_rect(s, x, CTOP + Inches(0.42) + i * CROW_H, w, CROW_H,
                 fill_color=bg, line_color=GRAY4, line_width=Pt(0.5))
        add_textbox(s, cell,
                    x + Inches(0.1),
                    CTOP + Inches(0.42) + i * CROW_H + Inches(0.08),
                    w - Inches(0.15), CROW_H - Inches(0.1),
                    font_size=Pt(10), color=GRAY2,
                    align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
                    bold=(i == 0 and j == 0))

add_textbox(s, "✓ = full support  |  ~ = partial  |  — = not available  |  ✗ = not applicable",
            Inches(0.5), Inches(6.65), Inches(12), Inches(0.35),
            font_size=Pt(9), color=GRAY3, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 35 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, GRAY6)
top_bar(s)
bottom_bar(s)
slide_number(s)
section_label(s, "05  Operations & Roadmap")
title_text(s, "What comes next", y=Inches(1.5), size=Pt(32))

TL_Y  = Inches(3.2)
TL_H  = Inches(0.08)
TL_X  = Inches(0.8)
TL_W  = Inches(11.8)
add_rect(s, TL_X, TL_Y, TL_W, TL_H, fill_color=GRAY3)

milestones = [
    (0.0,  "Today\nv1.18.0", "AI-action firewall + dynamic policy (recommend · grants · approve-and-learn)\nComposable by group · Multi-CA · Recording · Fail-closed + red-team hardening"),
    (0.28, "Near-term",      "Teams approval bridge (Entra)\nAWS KMS / GCP Cloud HSM\nControl-plane PKI cert · KRL"),
    (0.57, "Mid-term",       "Multi-instance sessions (Redis)\nWORM audit log export"),
    (0.85, "Long-term",      "Audit dashboard\nDynamic host registration"),
]

for frac, label, detail in milestones:
    x   = TL_X + frac * TL_W
    DOT = Inches(0.18)
    add_rect(s, x - DOT / 2, TL_Y - DOT / 2 + TL_H / 2,
             DOT, DOT, fill_color=BLACK)
    add_textbox(s, label,
                x - Inches(1.0), TL_Y - Inches(1.05),
                Inches(2.0), Inches(0.8),
                font_size=Pt(11), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_textbox(s, detail,
                x - Inches(1.2), TL_Y + Inches(0.35),
                Inches(2.4), Inches(1.2),
                font_size=Pt(10), color=GRAY2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 36 — SUMMARY / CTA
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
slide_bg(s, BLACK)
top_bar(s, bg=WHITE, height=Inches(0.08))
bottom_bar(s)
slide_number(s)

add_textbox(s, "SUMMARY",
            Inches(0.9), Inches(0.95), Inches(11), Inches(0.4),
            font_size=Pt(9), bold=True, color=GRAY4)

add_textbox(s, "SSH Broker gives AI agents\nsecure, audited infrastructure access —\nwithout ever handing them a key.",
            Inches(0.9), Inches(1.45), Inches(11.0), Inches(2.0),
            font_size=Pt(28), bold=True, color=WHITE)

takeaways = [
    "Ephemeral credentials — no static secrets, no exfiltration risk.",
    "Policy-gated by an isolated signer — broker compromise = minimal blast radius.",
    "Three-phase AI-action firewall — composable command policy, human approval, behaviour guardrails.",
    "Cryptographically chained audit trail — full attribution across signer + broker + sshd.",
    "Federated identity via Microsoft Entra ID — no new user directory to manage.",
    "Single Go binary per role — no cluster, no external dependencies required today.",
]

for i, t in enumerate(takeaways):
    add_rect(s, Inches(0.85), Inches(3.55) + i * Inches(0.57) + Inches(0.18),
             Inches(0.08), Inches(0.08), fill_color=WHITE)
    add_textbox(s, t,
                Inches(1.1), Inches(3.55) + i * Inches(0.57),
                Inches(11.5), Inches(0.55),
                font_size=Pt(13), color=GRAY4)

# ── Save ───────────────────────────────────────────────────────────────────
# Write next to this script so the generator is portable across machines.
OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "ssh_broker_presentation.pptx")
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")
